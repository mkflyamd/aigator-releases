"""OneDrive skill -- 2 tools."""

from pathlib import Path

ONEDRIVE_SKILLS_DIR = Path(__file__).parent.parent / "m365-onedrive" / "scripts"

SKILL_ID = "onedrive"
ALWAYS_ON = True


# A Graph drive-item id is an opaque token we can GET directly. We can't
# reliably recognize it by format — Microsoft uses multiple schemes
# (base32 "01ABC..." for personal OneDrive, base64url with '-'/'_' for
# SharePoint). Instead, is_resolvable_item_id() only filters out things
# that are DEFINITELY not an id (filenames, fallback markers). The actual
# "is this a real id?" check is a Graph GET, with name-search fallback.
import re as _re

_UNRESOLVABLE_RE = _re.compile(
    r"^onedrive:|^spo[@:]|\.(docx|pptx|xlsx|pdf|txt|md|csv|json|py|js|ts|html|xml|yaml|yml)$|[\s/]",
    _re.IGNORECASE,
)


def is_resolvable_item_id(item_id: str) -> bool:
    """True if item_id is NOT obviously a fallback/filename and is worth
    attempting a direct Graph GET. Never raises — unknown ids default True
    so the GET path decides. Use this only to skip obvious non-ids cheaply."""
    s = (item_id or "").strip()
    if not s:
        return False
    return not _UNRESOLVABLE_RE.search(s)


def _try_direct_lookup(
    gc, item_id: str, drive_id: str = "", select: str = "id,name,webUrl,parentReference"
) -> dict | None:
    """Attempt to resolve an item_id directly via Graph.

    Tries /drives/{drive_id}/items/{id} (or /me/drive/items/{id} when no
    drive_id), then falls back to /me/drive/sharedWithMe on 400/404 to discover
    the correct drive_id for SharePoint items pinned without one.

    Returns a dict with the raw Graph `meta` plus extracted fields on success,
    None if the id is not directly resolvable. Never raises — Graph is the
    source of truth for "is this a real id", not a regex.
    """
    if not item_id:
        return None

    def _from_meta(meta: dict, drive_id: str) -> dict:
        pr = meta.get("parentReference") or {}
        return {
            "id": meta.get("id", item_id),
            "drive_id": drive_id or pr.get("driveId", ""),
            "web_url": meta.get("webUrl", ""),
            "name": meta.get("name", ""),
            "meta": meta,
        }

    # 1. Direct GET — /drives/{drive_id}/items/{id} if drive_id known,
    #    else /me/drive/items/{id} (personal OneDrive).
    try:
        if drive_id:
            meta = gc.get(
                f"/drives/{drive_id}/items/{item_id}", params={"$select": select}
            )
        else:
            meta = gc.get(f"/me/drive/items/{item_id}", params={"$select": select})
        return _from_meta(meta, drive_id)
    except Exception as e:
        status = getattr(e, "status_code", 0)
        if status not in (400, 404):
            # Network/auth/throttle — give up on this path, don't mask the error
            # with name-search (which would also fail).
            return None
        # 400/404: id may be valid but on a SharePoint drive we don't know about.
        # Fall through to sharedWithMe discovery only when we have no drive_id.
        if drive_id:
            return None

    # 2. sharedWithMe fallback: scan for an item whose remoteItem.id matches,
    #    discover the real driveId, then GET via /drives/{driveId}/items/{id}.
    try:
        shared = gc.get(
            "/me/drive/sharedWithMe",
            params={
                "$select": "id,name,remoteItem",
                "$top": "200",
            },
        )
        for it in shared.get("value", []):
            ri = it.get("remoteItem") or {}
            ri_id = ri.get("id", "")
            if ri_id == item_id or it.get("id") == item_id:
                ri_drive = (ri.get("parentReference") or {}).get("driveId", "")
                if not ri_drive:
                    continue
                meta = gc.get(
                    f"/drives/{ri_drive}/items/{ri_id}", params={"$select": select}
                )
                return _from_meta(meta, ri_drive)
    except Exception:
        pass
    return None


def _try_direct_lookup_with_download_url(
    gc, item_id: str, drive_id: str = ""
) -> dict | None:
    """Resolve an item for download — includes @microsoft.graph.downloadUrl and size."""
    return _try_direct_lookup(
        gc,
        item_id,
        drive_id,
        select="id,name,size,file,webUrl,parentReference,@microsoft.graph.downloadUrl",
    )


def resolve_onedrive_item(
    filename: str, location_hint: str = "", item_id: str = "", drive_id: str = ""
) -> dict:
    """Resolve a OneDrive/SharePoint file to a real Graph drive-item.

    Resolution order (Graph is the source of truth — no regex allowlisting):
      1. If item_id is provided and not obviously a fallback, try a direct
         Graph GET (with sharedWithMe fallback on 404). This handles every
         id format Microsoft uses (base32 "01...", base64url SharePoint ids).
      2. Fall back to name-search via /search/query (driveItem), which spans
         OneDrive + SharePoint. Used when no id is available or the id failed
         direct lookup.

    Returns {"id","drive_id","web_url","name"} on a confident match, or
    {"error": ...} if nothing matched. Never guesses across ambiguity.
    """
    from .._m365.helpers import get_skill_client

    gc = get_skill_client(ONEDRIVE_SKILLS_DIR)

    # 1. Direct id lookup — try first when we have a plausible id.
    if item_id and is_resolvable_item_id(item_id):
        direct = _try_direct_lookup(gc, item_id, drive_id)
        if direct and direct.get("id"):
            # Strip the raw meta — callers only want the resolved fields.
            return {k: v for k, v in direct.items() if k != "meta"}

    clean = (filename or "").strip()
    if not clean and not item_id:
        return {"error": "empty filename and no item_id"}

    # 2. Name-search via the Microsoft Search API (spans OneDrive + SharePoint).
    if not clean:
        return {
            "error": f"could not resolve item_id {item_id!r} directly "
            f"and no filename provided for search"
        }
    # Sanitize the query: Microsoft Search treats '-' as a NOT operator, which
    # breaks queries like "PLANNING -ROCm" (returns everything BUT ROCm).
    # Replace hyphens/dashes with spaces and collapse whitespace.
    import re as _re_q

    query = _re_q.sub(r"[\-–—]+", " ", clean)
    query = _re_q.sub(r"\s+", " ", query).strip()
    try:
        body = {
            "requests": [
                {
                    "entityTypes": ["driveItem"],
                    "query": {"queryString": query},
                    "from": 0,
                    "size": 25,
                    "fields": ["name", "id", "parentReference", "webUrl"],
                }
            ]
        }
        data = gc.post("/search/query", body)
    except Exception as e:
        return {"error": f"search failed: {e}"}

    hits = []
    for resp in data.get("value", []):
        for container in resp.get("hitsContainers", []):
            for hit in container.get("hits", []):
                res = hit.get("resource", {}) or {}
                hits.append(res)

    if not hits:
        return {"error": f"no file matched {clean!r}"}

    target = clean.lower()

    def _hit_fields(res):
        name = res.get("name", "")
        rid = res.get("id", "")
        pr = res.get("parentReference", {}) or {}
        d_id = pr.get("driveId", "")
        web_url = res.get("webUrl", "")
        return name, rid, d_id, web_url

    # Exact filename match, optionally disambiguated by location_hint.
    exact = [h for h in hits if _hit_fields(h)[0].lower() == target]
    loc = (location_hint or "").strip().lower()
    if exact and loc:
        located = [h for h in exact if loc in (_hit_fields(h)[3] or "").lower()]
        if located:
            exact = located
    if exact:
        name, rid, d_id, web_url = _hit_fields(exact[0])
        return {"id": rid, "drive_id": d_id, "web_url": web_url, "name": name}

    # Unique prefix match (filename without extension) — only if unambiguous.
    stem = target.rsplit(".", 1)[0]
    prefix = [h for h in hits if _hit_fields(h)[0].lower().startswith(stem)]
    if len(prefix) == 1:
        name, rid, d_id, web_url = _hit_fields(prefix[0])
        return {"id": rid, "drive_id": d_id, "web_url": web_url, "name": name}

    # Fuzzy normalized match — handles DOM-scraped labels that lost a zero-width
    # space, non-breaking space, missing extension, or extra whitespace vs. the
    # real filename. Normalization: strip zero-width + non-breaking spaces,
    # replace non-alphanumerics with spaces, collapse whitespace, strip a known
    # file extension if present, lowercase.
    import re as _re_norm

    _KNOWN_EXTS = (
        ".docx",
        ".pptx",
        ".xlsx",
        ".pdf",
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".py",
        ".js",
        ".ts",
        ".html",
        ".xml",
        ".yaml",
        ".yml",
        ".mp4",
        ".mov",
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
    )

    def _norm(s: str) -> str:
        s = s.replace("\u200b", "").replace("\u200c", "").replace("\u200d", "")
        s = s.replace("\ufeff", "").replace("\xa0", " ")
        low = s.lower()
        for ext in _KNOWN_EXTS:
            if low.endswith(ext):
                s = s[: -len(ext)]
                break
        s = _re_norm.sub(r"[^a-z0-9]+", " ", s.lower())  # non-alnum -> single space
        s = _re_norm.sub(r"\s+", " ", s).strip()
        return s

    norm_target = _norm(clean)
    if norm_target:
        norm_hits = [h for h in hits if _norm(_hit_fields(h)[0]) == norm_target]
        if norm_hits and loc:
            located = [h for h in norm_hits if loc in (_hit_fields(h)[3] or "").lower()]
            if located:
                norm_hits = located
        if len(norm_hits) == 1:
            name, rid, d_id, web_url = _hit_fields(norm_hits[0])
            return {"id": rid, "drive_id": d_id, "web_url": web_url, "name": name}

    return {
        "error": f"no confident match for {clean!r} "
        f"({len(hits)} candidates, none exact)"
    }


TOOL_DEFS = [
    {
        "name": "read_onedrive_file",
        "description": (
            "Download and read the text content of a file from OneDrive or SharePoint. "
            "Use when user asks to open, read, summarize, or extract content from a specific file. "
            "Supports .docx, .txt, .md, .csv, .xlsx, .pptx, and plain text formats. "
            "Provide file_id (and drive_id for shared/SharePoint files), a file_path, or a SharePoint share-link URL. "
            "For large .pptx decks use slide_start/slide_end to read a specific range of slides (1-based, inclusive). "
            "Call get_pptx_info first to learn the total slide count, then page through in batches."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "file_id": {
                    "type": "string",
                    "description": "OneDrive item ID (from a prior search/list or pinned file)",
                },
                "drive_id": {
                    "type": "string",
                    "description": "Drive ID for files shared with you or hosted on SharePoint. Required when the file is not on the user's own OneDrive.",
                },
                "file_path": {
                    "type": "string",
                    "description": "File path relative to OneDrive root, e.g. 'Documents/report.docx'",
                },
                "share_url": {
                    "type": "string",
                    "description": "A SharePoint share-link URL (e.g. https://tenant-my.sharepoint.com/:w:/g/personal/...). Gator will resolve it to a real file via Graph.",
                },
                "filename_hint": {
                    "type": "string",
                    "description": "The file's display name/label (e.g. from a pinned context). Used as a fallback to search by name when file_id cannot be resolved directly. Pass this when you have a label but aren't sure the file_id is a valid Graph item ID.",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "Max characters to return. Default 200000 for .pptx, 8000 for other types.",
                    "default": 8000,
                },
                "slide_start": {
                    "type": "integer",
                    "description": "For .pptx: 1-based first slide to include. Default 1 (start of deck).",
                },
                "slide_end": {
                    "type": "integer",
                    "description": "For .pptx: 1-based last slide to include (inclusive). Default: all slides.",
                },
            },
            "required": [],
        },
    },
    {
        "name": "list_onedrive_files",
        "description": "List files and folders in the user's OneDrive. Use when user asks about their OneDrive, files, or documents.",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Folder path to list (e.g. 'Documents/Projects'). Default: root.",
                    "default": "",
                },
                "count": {
                    "type": "integer",
                    "description": "Max items. Default 50.",
                    "default": 50,
                },
            },
            "required": [],
        },
    },
    {
        "name": "search_onedrive_files",
        "description": "Search for files in OneDrive by name or content. Use when user asks to find a file.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query"},
                "count": {
                    "type": "integer",
                    "description": "Max results. Default 10.",
                    "default": 10,
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "download_onedrive_file",
        "description": (
            "Download a file from OneDrive or SharePoint and save it to the local disk as raw bytes. "
            "Use when you need the actual file (e.g. to read hyperlinks from a .docx, edit a .pptx with python-pptx, "
            "or work with the file locally). Returns the local path where the file was saved. "
            "Provide file_id (and drive_id for SharePoint files), file_path, or share_url. "
            "Optionally provide local_path to control where the file is saved (default: ~/Downloads/<filename>)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "file_id": {"type": "string", "description": "OneDrive item ID"},
                "drive_id": {
                    "type": "string",
                    "description": "Drive ID for SharePoint files",
                },
                "file_path": {
                    "type": "string",
                    "description": "File path relative to OneDrive root",
                },
                "share_url": {
                    "type": "string",
                    "description": "SharePoint share-link URL",
                },
                "local_path": {
                    "type": "string",
                    "description": "Where to save the file locally. Default: ~/Downloads/<filename>",
                },
            },
            "required": [],
        },
    },
]

TOOL_STATUS = {
    "read_onedrive_file": "\U0001f4c4 Reading file...",
    "list_onedrive_files": "\U0001f4c1 Browsing OneDrive...",
    "search_onedrive_files": "\U0001f50d Searching OneDrive...",
    "download_onedrive_file": "\U0001f4e5 Downloading file...",
}


def _tool_read_onedrive_file(
    file_id: str = "",
    drive_id: str = "",
    file_path: str = "",
    share_url: str = "",
    max_chars: int = 0,
    slide_start: int = 1,
    slide_end: int = 0,
    filename_hint: str = "",
    _context_id: str = "",
) -> dict:
    import io
    from .._m365.helpers import get_skill_client

    gc = get_skill_client(ONEDRIVE_SKILLS_DIR)

    if not file_id and not file_path and not share_url:
        return {"error": "Provide file_id, file_path, or share_url"}

    # Resolve SharePoint share-links via Graph /shares/{encodedUrl}/driveItem.
    # This is the only correct way — never fall back to a name search.
    if share_url:
        import base64

        encoded = (
            base64.b64encode(share_url.encode())
            .decode()
            .rstrip("=")
            .replace("/", "_")
            .replace("+", "-")
        )
        share_token = f"u!{encoded}"
        try:
            resolved = gc.get(
                f"/shares/{share_token}/driveItem",
                params={"$select": "id,name,parentReference"},
            )
            file_id = resolved["id"]
            drive_id = resolved.get("parentReference", {}).get("driveId", drive_id)
        except Exception as e:
            return {
                "error": (
                    f"Could not resolve SharePoint share-link: {e}. "
                    "Please confirm the filename or paste the file content directly."
                ),
                "share_url": share_url,
                "unresolved": True,
            }

    # Resolve item metadata — include @microsoft.graph.downloadUrl upfront.
    # This pre-authenticated URL is the most reliable way to download SharePoint-hosted
    # files (MySite, Teams files, etc.) because it bypasses the auth redirect chain.
    #
    # Resolution strategy (Graph is the source of truth — no regex allowlisting):
    #   1. file_path → /me/drive/root:/<path> (direct path lookup)
    #   2. file_id → _try_direct_lookup (direct GET, sharedWithMe fallback on
    #      400/404 for SharePoint items pinned without a drive_id)
    #   3. name-search via resolve_onedrive_item (last resort — used when the
    #      id is a fallback marker like "onedrive:filename")
    meta = None
    if file_path and not file_id:
        from urllib.parse import quote

        meta_path = f"/me/drive/root:/{quote(file_path.lstrip('/'))}"
        meta = gc.get(
            meta_path,
            params={
                "$select": "id,name,size,file,webUrl,parentReference,@microsoft.graph.downloadUrl"
            },
        )
    elif file_id:
        direct = _try_direct_lookup_with_download_url(gc, file_id, drive_id)
        if direct:
            meta = direct["meta"]
            file_id = meta.get("id", file_id)
            drive_id = direct["drive_id"]

    if meta is None:
        # Direct lookup failed — try the pin's web_url as a share link next.
        # The shell captures a share-link URL for SharePoint files that Graph can
        # resolve via /shares/{token}/driveItem. This is more reliable than
        # name-search (no ambiguity) and works for any SharePoint file.
        if _context_id and file_id:
            try:
                from skills.context.state import get_pins

                for pin in get_pins(_context_id):
                    if pin.get("source") == "onedrive" and pin.get("id") == file_id:
                        pin_web_url = (pin.get("meta") or {}).get("web_url", "")
                        if (
                            pin_web_url
                            and "sharepoint.com" in pin_web_url
                            and "/:" in pin_web_url
                        ):
                            # Looks like a SharePoint share link — resolve it.
                            import base64

                            encoded = (
                                base64.b64encode(pin_web_url.encode())
                                .decode()
                                .rstrip("=")
                                .replace("/", "_")
                                .replace("+", "-")
                            )
                            share_token = f"u!{encoded}"
                            try:
                                resolved = gc.get(
                                    f"/shares/{share_token}/driveItem",
                                    params={
                                        "$select": "id,name,size,file,webUrl,parentReference,@microsoft.graph.downloadUrl"
                                    },
                                )
                                meta = resolved
                                file_id = resolved.get("id", file_id)
                                drive_id = resolved.get("parentReference", {}).get(
                                    "driveId", drive_id
                                )
                            except Exception:
                                pass  # fall through to name-search
                        break
            except Exception:
                pass

    if meta is None:
        # Direct lookup AND share-link both failed — last resort: name-search
        # via the Microsoft Search API. Derive a filename from file_path,
        # filename_hint, or file_id (which may be a "onedrive:filename" marker).
        filename = (file_path or filename_hint or "").strip()
        # If no explicit hint, look up the pin's label from the tab context.
        # The server injects _context_id — this lets read_onedrive_file resolve
        # a bad pin id automatically, without relying on the agent to pass
        # filename_hint.
        if not filename and _context_id and file_id:
            try:
                from skills.context.state import get_pins

                for pin in get_pins(_context_id):
                    if pin.get("source") == "onedrive" and pin.get("id") == file_id:
                        filename = (pin.get("label") or "").strip()
                        break
            except Exception:
                pass
        if not filename and file_id and ":" in file_id:
            filename = file_id.split(":", 1)[1]
        if filename:
            resolved = resolve_onedrive_item(
                filename=filename, item_id=file_id, drive_id=drive_id
            )
            if not resolved.get("error") and resolved.get("id"):
                file_id = resolved["id"]
                drive_id = resolved.get("drive_id", drive_id)
                if drive_id:
                    meta = gc.get(
                        f"/drives/{drive_id}/items/{file_id}",
                        params={
                            "$select": "id,name,size,file,webUrl,parentReference,@microsoft.graph.downloadUrl"
                        },
                    )
                else:
                    meta = gc.get(
                        f"/me/drive/items/{file_id}",
                        params={
                            "$select": "id,name,size,file,webUrl,parentReference,@microsoft.graph.downloadUrl"
                        },
                    )
        if meta is None:
            return {
                "error": (
                    f"Could not resolve file_id {file_id!r} (drive_id={drive_id!r}). "
                    "The file may be on a SharePoint site you can access — try sharing "
                    "the file link, or tell me the exact filename and I'll search for it."
                ),
                "file_id": file_id,
                "drive_id": drive_id,
            }
    item_id = meta.get("id", file_id)
    name = meta.get("name", "")
    web_url = meta.get("webUrl", "")
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    direct_url = meta.get("@microsoft.graph.downloadUrl", "")

    import httpx

    token = gc.get_token()

    # Module-level pool for file downloads (reuses TCP connections)
    global _od_dl_pool
    if (
        not hasattr(_tool_read_onedrive_file, "_pool")
        or _tool_read_onedrive_file._pool.is_closed
    ):
        _tool_read_onedrive_file._pool = httpx.Client(
            timeout=httpx.Timeout(60.0), follow_redirects=True
        )
    _pool = _tool_read_onedrive_file._pool

    def _download() -> bytes:
        # Prefer the pre-authenticated downloadUrl (no Authorization header needed).
        if direct_url:
            r = _pool.get(direct_url)
            r.raise_for_status()
            return r.content
        # Fall back to the Graph /content endpoint with Bearer token.
        if drive_id:
            dl_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}/content"
        else:
            dl_url = (
                f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/content"
            )
        r = _pool.get(dl_url, headers={"Authorization": f"Bearer {token}"})
        r.raise_for_status()
        return r.content

    raw = _download()

    # Guard: if response looks like HTML (e.g. a login redirect that returned 200),
    # the token is wrong or the file requires different permissions.
    if raw[:5] in (b"<!DOC", b"<html", b"<HTML") or raw[:3] == b"\xef\xbb\xbf<":
        return {
            "error": "OneDrive returned an HTML page instead of the file — the token may have expired or lack Files.Read scope. Try refreshing your OneDrive token in Settings.",
            "name": name,
            "url": web_url,
            "auth_required": True,
        }

    def _docx_extract_xml(data: bytes) -> str:
        """
        Read text directly from the docx ZIP without using python-docx.
        A .docx is a ZIP containing word/document.xml — parse w:t elements directly.
        This tolerates the strict-OOXML / Word Online save format that python-docx rejects.
        """
        import zipfile
        import xml.etree.ElementTree as ET

        W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        P = f"{{{W}}}p"
        TBL = f"{{{W}}}tbl"
        TR = f"{{{W}}}tr"
        TC = f"{{{W}}}tc"
        T = f"{{{W}}}t"
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                with z.open("word/document.xml") as f:
                    root = ET.parse(f).getroot()
            body = root.find(f"{{{W}}}body")
            if body is None:
                return ""
            lines = []
            for child in body:
                if child.tag == P:
                    line = "".join(t.text for t in child.iter(T) if t.text).strip()
                    if line:
                        lines.append(line)
                elif child.tag == TBL:
                    rows = []
                    for tr in child.iter(TR):
                        cells = [
                            "".join(t.text for t in tc.iter(T) if t.text).strip()
                            for tc in tr.findall(f"{{{W}}}tc")
                        ]
                        if cells:
                            rows.append(cells)
                    if rows:
                        col_count = max(len(r) for r in rows)
                        header = "| " + " | ".join(rows[0]) + " |"
                        sep = "| " + " | ".join(["---"] * col_count) + " |"
                        body_rows = [
                            "| " + " | ".join(r + [""] * (col_count - len(r))) + " |"
                            for r in rows[1:]
                        ]
                        lines.append("\n".join([header, sep] + body_rows))
            return "\n".join(lines)
        except Exception:
            return ""

    def _docx_extract_hyperlinks(data: bytes) -> list[dict]:
        """Extract hyperlinks from a .docx: display text + target URL pairs.

        Reads word/_rels/document.xml.rels for relationship targets, then
        walks word/document.xml to find <w:hyperlink r:id="..."> elements and
        collects the display text of each linked run.
        Returns a list of {text, url} dicts (deduped, ordered by appearance).
        """
        import zipfile
        import xml.etree.ElementTree as ET

        W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        HL = f"{{{W}}}hyperlink"
        T = f"{{{W}}}t"
        RID = f"{{{R}}}id"
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                # Build rid → url map from relationships file
                rels: dict[str, str] = {}
                if "word/_rels/document.xml.rels" in z.namelist():
                    with z.open("word/_rels/document.xml.rels") as f:
                        rels_root = ET.parse(f).getroot()
                    for rel in rels_root:
                        rid = rel.attrib.get("Id", "")
                        target = rel.attrib.get("Target", "")
                        typ = rel.attrib.get("Type", "")
                        if "hyperlink" in typ and target:
                            rels[rid] = target
                if not rels:
                    return []
                # Walk document.xml for <w:hyperlink r:id="..."> elements
                with z.open("word/document.xml") as f:
                    doc_root = ET.parse(f).getroot()
            seen: dict[str, str] = {}  # url → first display text
            results: list[dict] = []
            for hl in doc_root.iter(HL):
                rid = hl.attrib.get(RID, "")
                if not rid or rid not in rels:
                    continue
                url = rels[rid]
                text = "".join(t.text for t in hl.iter(T) if t.text).strip()
                if not text:
                    text = url
                if url not in seen:
                    seen[url] = text
                    results.append({"text": text, "url": url})
            return results
        except Exception:
            return []

    def _docx_extract_images(data: bytes, max_images: int = 5) -> list[dict]:
        """Extract embedded images from a docx ZIP (word/media/*)."""
        import zipfile
        import base64

        images = []
        _IMG_EXTS = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".bmp": "image/bmp",
            ".tiff": "image/tiff",
        }
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                media_files = sorted(
                    [n for n in z.namelist() if n.startswith("word/media/")]
                )
                for name in media_files[:max_images]:
                    ext = "." + name.rsplit(".", 1)[-1].lower() if "." in name else ""
                    mime = _IMG_EXTS.get(ext)
                    if not mime:
                        continue
                    img_data = z.read(name)
                    if len(img_data) < 500:
                        continue  # skip tiny images (icons, bullets)
                    images.append(
                        {
                            "name": name.split("/")[-1],
                            "media_type": mime,
                            "base64": base64.b64encode(img_data).decode("ascii"),
                        }
                    )
        except Exception:
            pass
        return images

    # Extract text based on file type
    text = ""
    if ext == "docx":
        try:
            import docx

            doc = docx.Document(io.BytesIO(raw))
            parts = []
            # Walk body children in document order so tables appear inline with paragraphs
            for block in doc.element.body:
                tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag
                if tag == "p":
                    from docx.oxml.ns import qn

                    text_parts = [t.text for t in block.iter(qn("w:t")) if t.text]
                    line = "".join(text_parts).strip()
                    if line:
                        parts.append(line)
                elif tag == "tbl":
                    from docx.oxml.ns import qn

                    rows = []
                    for tr in block.iter(qn("w:tr")):
                        cells = []
                        for tc in tr.iter(qn("w:tc")):
                            cell_parts = [t.text for t in tc.iter(qn("w:t")) if t.text]
                            cells.append("".join(cell_parts).strip())
                        if cells:
                            rows.append(cells)
                    if rows:
                        col_count = max(len(r) for r in rows)
                        header = "| " + " | ".join(rows[0]) + " |"
                        sep = "| " + " | ".join(["---"] * col_count) + " |"
                        body_rows = [
                            "| " + " | ".join(r + [""] * (col_count - len(r))) + " |"
                            for r in rows[1:]
                        ]
                        parts.append("\n".join([header, sep] + body_rows))
            text = "\n".join(parts)
        except Exception:
            # python-docx uses a strict ZIP parser and rejects Word Online's extended format.
            # If the file is OLE2 (AIP/IRM-protected), fall back to server-side PDF conversion.
            if raw[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
                return {
                    "error": (
                        "This file has a sensitivity label (AIP/IRM) that encrypts it. "
                        "To read it: open the file in Word desktop or Word Online → click the "
                        "Sensitivity button in the ribbon → change the label to General or remove it → save. "
                        "Then ask me to read it again."
                    ),
                    "name": name,
                    "url": web_url,
                }
            else:
                # Regular ZIP-based docx that python-docx rejected — try raw XML.
                text = _docx_extract_xml(raw)
                if not text:
                    return {
                        "error": (
                            "Could not parse the .docx file locally. "
                            "The file may be password-protected. "
                            "Open it in Word Online and paste the text here."
                        ),
                        "name": name,
                        "url": web_url,
                    }
    elif ext in ("xlsx", "xls"):
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                if any(c is not None for c in row):
                    parts.append("\t".join("" if c is None else str(c) for c in row))
        text = "\n".join(parts)
    elif ext in ("pptx",):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw))
        total_slides = len(prs.slides)
        _s_start = max(1, slide_start or 1)
        _s_end = (
            min(total_slides, slide_end)
            if slide_end
            else min(total_slides, _s_start + 199)
        )
        parts = [f"Total slides: {total_slides}  |  Showing slides {_s_start}–{_s_end}"]
        for i, slide in enumerate(prs.slides, 1):
            if i < _s_start or i > _s_end:
                continue
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        text = "\n".join(parts)
        if not max_chars:
            max_chars = 200000
    elif ext in (
        "txt",
        "md",
        "csv",
        "json",
        "py",
        "js",
        "ts",
        "html",
        "xml",
        "yaml",
        "yml",
    ):
        text = raw.decode("utf-8", errors="replace")
    else:
        # Attempt UTF-8 decode for unknown types
        try:
            text = raw.decode("utf-8", errors="strict")
        except Exception:
            return {
                "error": f"Cannot extract text from .{ext} files",
                "name": name,
                "url": web_url,
            }

    if not max_chars:
        max_chars = 8000
    truncated = len(text) > max_chars
    result = {
        "name": name,
        "url": web_url,
        "size_bytes": len(raw),
        "truncated": truncated,
        "content": text[:max_chars] + ("\n\n[... truncated ...]" if truncated else ""),
    }
    if ext == "pptx":
        result["total_slides"] = total_slides
        result["slides_shown"] = f"{_s_start}–{_s_end}"
    # Extract embedded images and hyperlinks from docx
    if ext == "docx":
        images = _docx_extract_images(raw)
        print(
            f"[onedrive] docx image extraction: found {len(images)} images in {name}",
            flush=True,
        )
        if images:
            result["_images"] = images
            result["_images_found"] = len(images)
        hyperlinks = _docx_extract_hyperlinks(raw)
        if hyperlinks:
            result["hyperlinks"] = hyperlinks
            print(
                f"[onedrive] docx hyperlink extraction: found {len(hyperlinks)} links in {name}",
                flush=True,
            )
    return result


def _tool_list_onedrive_files(path: str = "", count: int = 50) -> dict:
    from .._m365.helpers import get_skill_client

    gc = get_skill_client(ONEDRIVE_SKILLS_DIR)
    api_path = (
        f"/me/drive/root:/{path}:/children" if path else "/me/drive/root/children"
    )
    data = gc.get(
        api_path,
        params={
            "$top": str(count),
            "$orderby": "name",
            "$select": "name,size,lastModifiedDateTime,folder,file,webUrl,id",
        },
    )
    items = []
    for item in data.get("value", []):
        is_folder = "folder" in item
        items.append(
            {
                "name": item.get("name", ""),
                "type": "folder" if is_folder else "file",
                "size": item.get("size", 0),
                "modified": item.get("lastModifiedDateTime", "")[:16],
                "url": item.get("webUrl", ""),
                "id": item.get("id", ""),
            }
        )
    return {"path": path or "/", "total": len(items), "items": items}


def _tool_search_onedrive_files(query: str, count: int = 10) -> dict:
    from .._m365.helpers import get_skill_client

    gc = get_skill_client(ONEDRIVE_SKILLS_DIR)
    data = gc.get(
        f"/me/drive/root/search(q='{query}')",
        params={
            "$top": str(count),
            "$select": "name,size,lastModifiedDateTime,parentReference,webUrl,id",
        },
    )
    items = []
    for item in data.get("value", []):
        parent_path = (
            item.get("parentReference", {})
            .get("path", "")
            .replace("/drive/root:", "")
            .lstrip("/")
        )
        items.append(
            {
                "name": item.get("name", ""),
                "path": f"{parent_path}/{item.get('name', '')}"
                if parent_path
                else item.get("name", ""),
                "size": item.get("size", 0),
                "modified": item.get("lastModifiedDateTime", "")[:16],
                "url": item.get("webUrl", ""),
                "id": item.get("id", ""),
                "drive_id": item.get("parentReference", {}).get("driveId", ""),
            }
        )

    # Personal-drive search misses SharePoint/shared files. If it came back empty,
    # fall back to the Microsoft Search API, which spans OneDrive + SharePoint.
    if not items:
        try:
            body = {
                "requests": [
                    {
                        "entityTypes": ["driveItem"],
                        "query": {"queryString": query},
                        "from": 0,
                        "size": count,
                        "fields": [
                            "name",
                            "id",
                            "parentReference",
                            "webUrl",
                            "size",
                            "lastModifiedDateTime",
                        ],
                    }
                ]
            }
            sdata = gc.post("/search/query", body)
            for resp in sdata.get("value", []):
                for container in resp.get("hitsContainers", []):
                    for hit in container.get("hits", []):
                        res = hit.get("resource", {}) or {}
                        pr = res.get("parentReference", {}) or {}
                        items.append(
                            {
                                "name": res.get("name", ""),
                                "path": res.get("name", ""),
                                "size": res.get("size", 0),
                                "modified": (res.get("lastModifiedDateTime", "") or "")[
                                    :16
                                ],
                                "url": res.get("webUrl", ""),
                                "id": res.get("id", ""),
                                "drive_id": pr.get("driveId", ""),
                            }
                        )
        except Exception as e:
            return {"query": query, "total": 0, "items": [], "search_api_error": str(e)}

    return {"query": query, "total": len(items), "items": items}


def _tool_download_onedrive_file(
    file_id: str = "",
    drive_id: str = "",
    file_path: str = "",
    share_url: str = "",
    local_path: str = "",
) -> dict:
    """Download a OneDrive/SharePoint file and save it to disk as raw bytes.
    Returns the local path where the file was saved."""
    import io as _io
    from .._m365.helpers import get_skill_client
    import httpx as _httpx

    gc = get_skill_client(ONEDRIVE_SKILLS_DIR)

    if not file_id and not file_path and not share_url:
        return {"error": "Provide file_id, file_path, or share_url"}

    # Resolve share URL
    if share_url:
        import base64

        encoded = (
            base64.b64encode(share_url.encode())
            .decode()
            .rstrip("=")
            .replace("/", "_")
            .replace("+", "-")
        )
        share_token = f"u!{encoded}"
        try:
            resolved = gc.get(
                f"/shares/{share_token}/driveItem",
                params={"$select": "id,name,parentReference"},
            )
            file_id = resolved["id"]
            drive_id = resolved.get("parentReference", {}).get("driveId", drive_id)
        except Exception as e:
            return {"error": f"Could not resolve share URL: {e}"}

    # Get metadata + pre-authenticated download URL
    if file_id:
        meta_path = (
            f"/drives/{drive_id}/items/{file_id}"
            if drive_id
            else f"/me/drive/items/{file_id}"
        )
    else:
        from urllib.parse import quote

        meta_path = f"/me/drive/root:/{quote(file_path.lstrip('/'))}"
    meta = gc.get(
        meta_path, params={"$select": "id,name,size,@microsoft.graph.downloadUrl"}
    )
    name = meta.get("name", "file")
    direct_url = meta.get("@microsoft.graph.downloadUrl", "")
    item_id = meta.get("id", file_id)

    # Download bytes
    pool = _httpx.Client(timeout=_httpx.Timeout(120.0), follow_redirects=True)
    try:
        if direct_url:
            r = pool.get(direct_url)
        else:
            dl_url = (
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}/content"
                if drive_id
                else f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}/content"
            )
            r = pool.get(dl_url, headers={"Authorization": f"Bearer {gc.get_token()}"})
        r.raise_for_status()
        raw = r.content
    finally:
        pool.close()

    if raw[:5] in (b"<!DOC", b"<html", b"<HTML"):
        return {"error": "Got HTML instead of file bytes — token may have expired"}

    # Determine save path
    from pathlib import Path as _Path

    if local_path:
        dest = _Path(local_path).expanduser()
    else:
        dest = _Path.home() / "Downloads" / name
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(raw)
    return {"saved_to": str(dest), "name": name, "size_bytes": len(raw)}


TOOL_HANDLERS = {
    "read_onedrive_file": _tool_read_onedrive_file,
    "list_onedrive_files": _tool_list_onedrive_files,
    "search_onedrive_files": _tool_search_onedrive_files,
    "download_onedrive_file": _tool_download_onedrive_file,
}
