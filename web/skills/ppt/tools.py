"""PowerPoint skill -- 4 tools."""

import os
import contextlib
import tempfile

SKILL_ID = "ppt"
SKILL_ALIASES = ["ppt_skill"]
# Always-on: PowerPoint read/edit tools (read_pptx, update_pptx, pptx_apply_theme,
# table/shape tools) are cheap to expose and the model should be able to act on a
# .pptx the user names without first selecting the skill. SKILL.md prompt is still
# gated (injected only when ppt is auto/active) so no per-turn token bloat.
ALWAYS_ON = True

# ── Layout name mapping (python-pptx built-in layout indices) ────────────────
LAYOUT_MAP = {
    "title_slide": 0,
    "title_content": 1,
    "section": 2,
    "two_content": 3,
    "comparison": 4,
    "title_only": 5,
    "blank": 6,
}

TOOL_DEFS = [
    {
        "name": "get_pptx_info",
        "description": "Get structural info about a PowerPoint presentation: slide count, slide titles, layouts, and dimensions. Use this first to understand a presentation before reading or editing. Use file_path='open' for the currently open PowerPoint via COM. OneDrive item IDs are automatically downloaded, edited, and re-uploaded.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to .pptx file, OneDrive item ID (from search_onedrive_files or list_onedrive_files), or 'open' / 'open:Deck.pptx' for the active presentation via COM. OneDrive item IDs are automatically downloaded, edited, and re-uploaded.",
                },
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "read_pptx",
        "description": "Read text content from PowerPoint slides. Returns slide titles, body text, shape text, and speaker notes. Use file_path='open' for COM, or provide a full file path. Optionally read a single slide by number.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to .pptx file, or 'open' for the active presentation via COM, or 'open:Deck.pptx' for a specific open presentation",
                },
                "slide_number": {
                    "type": "integer",
                    "description": "1-based slide number to read. Omit to read all slides.",
                },
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "create_pptx",
        "description": "Create a new PowerPoint presentation from scratch. Provide an array of slide definitions with layout, title, content (bullets or text), and optional notes/images. Call ONCE with all slides.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path where the .pptx file will be saved",
                },
                "slides": {
                    "type": "array",
                    "description": "Array of slide definitions.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "layout": {
                                "type": "string",
                                "enum": [
                                    "title_slide",
                                    "title_content",
                                    "section",
                                    "blank",
                                    "two_content",
                                    "comparison",
                                    "title_only",
                                ],
                                "description": "Slide layout. Defaults to 'title_content'.",
                                "default": "title_content",
                            },
                            "title": {"type": "string", "description": "Slide title"},
                            "subtitle": {
                                "type": "string",
                                "description": "Subtitle (for title_slide layout only)",
                            },
                            "content": {
                                "description": "Slide body content. Array of strings for bullet points, or a single string for one paragraph.",
                                "oneOf": [
                                    {"type": "string"},
                                    {"type": "array", "items": {"type": "string"}},
                                ],
                            },
                            "notes": {
                                "type": "string",
                                "description": "Speaker notes for this slide",
                            },
                            "image": {
                                "type": "object",
                                "description": "Optional image to add to the slide.",
                                "properties": {
                                    "path": {
                                        "type": "string",
                                        "description": "Full path to image file",
                                    },
                                    "left": {
                                        "type": "number",
                                        "description": "Left position in inches. Default 1.",
                                        "default": 1,
                                    },
                                    "top": {
                                        "type": "number",
                                        "description": "Top position in inches. Default 2.",
                                        "default": 2,
                                    },
                                    "width": {
                                        "type": "number",
                                        "description": "Width in inches. Default 5.",
                                        "default": 5,
                                    },
                                },
                                "required": ["path"],
                            },
                        },
                    },
                },
                "author": {
                    "type": "string",
                    "description": "Presentation author metadata. Optional.",
                    "default": "",
                },
            },
            "required": ["file_path", "slides"],
        },
    },
    {
        "name": "update_pptx",
        "description": "Update PowerPoint slides. Supports single slide updates or batch mode for multiple slides in one call. ALWAYS use batch mode when updating more than one slide. Use file_path='open' for COM, or provide a full file path.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to .pptx file, or 'open' for the active presentation via COM, or 'open:Deck.pptx' for a specific open presentation. Defaults to 'open'.",
                    "default": "open",
                },
                "slide_number": {
                    "type": "integer",
                    "description": "1-based slide number to update (for single operation)",
                },
                "update_type": {
                    "type": "string",
                    "enum": ["title", "body", "shape", "batch"],
                    "description": "What to update. Use 'batch' to update multiple slides in one call.",
                },
                "new_text": {
                    "type": "string",
                    "description": "New text content to set (for single operation)",
                },
                "shape_index": {
                    "type": "integer",
                    "description": "0-based shape index when update_type is 'shape'. Defaults to 0.",
                    "default": 0,
                },
                "operations": {
                    "type": "array",
                    "description": "For batch mode: array of slide updates.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "slide_number": {"type": "integer"},
                            "update_type": {
                                "type": "string",
                                "enum": ["title", "body", "shape"],
                            },
                            "new_text": {"type": "string"},
                            "shape_index": {"type": "integer", "default": 0},
                        },
                        "required": ["slide_number", "update_type", "new_text"],
                    },
                },
            },
            "required": ["update_type"],
        },
    },
    {
        "name": "pptx_list_shapes",
        "description": "List every shape on a slide with its type (TABLE/PICTURE/TEXT_BOX/AUTO_SHAPE/GROUP), name, geometry (left/top/width/height in inches), whether it holds a table, and a text preview. Use this to discover shapes before editing. slide_locator may be a 1-based index OR a content string to scan for.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to the .pptx file (closed file; not COM)",
                },
                "slide_locator": {
                    "description": "1-based slide index, or a text string to find the slide by content"
                },
            },
            "required": ["file_path", "slide_locator"],
        },
    },
    {
        "name": "pptx_read_table",
        "description": "Read a table as a rows x cols grid. Each cell reports {text, fill_hex, font_hex, font_bold}. table_locator may be a 0-based table index OR a header-row cell text to scan for. Use before writing cells so you know the layout.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to the .pptx file",
                },
                "slide_locator": {
                    "description": "1-based slide index, or content string"
                },
                "table_locator": {
                    "description": "0-based table index, or a header-row cell text to match"
                },
            },
            "required": ["file_path", "slide_locator", "table_locator"],
        },
    },
    {
        "name": "pptx_write_table_cell",
        "description": "Write any subset of {text, fill_hex, font_hex, bold} into one table cell by row/col. Omitted fields are left untouched (partial update preserves existing text). Returns cell_after read back from the re-saved file so success is proven. Colors are 6-digit hex with no '#'.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to the .pptx file",
                },
                "slide_locator": {
                    "description": "1-based slide index, or content string"
                },
                "table_locator": {
                    "description": "0-based table index, or header-row cell text"
                },
                "row": {"type": "integer", "description": "0-based row index"},
                "col": {"type": "integer", "description": "0-based column index"},
                "text": {
                    "type": "string",
                    "description": "New cell text. Omit to leave text unchanged.",
                },
                "fill_hex": {
                    "type": "string",
                    "description": "Cell fill color, 6-digit hex no '#'. Omit to leave unchanged.",
                },
                "font_hex": {
                    "type": "string",
                    "description": "Font color, 6-digit hex no '#'. Omit to leave unchanged.",
                },
                "bold": {
                    "type": "boolean",
                    "description": "Font bold. Omit to leave unchanged.",
                },
            },
            "required": ["file_path", "slide_locator", "table_locator", "row", "col"],
        },
    },
    {
        "name": "pptx_add_table_row",
        "description": "Append a row to a table by deep-copying the last row's formatting. With copy_last=False the new cells are blanked. Returns the new row index and total row count read back from disk.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to the .pptx file",
                },
                "slide_locator": {
                    "description": "1-based slide index, or content string"
                },
                "table_locator": {
                    "description": "0-based table index, or header-row cell text"
                },
                "copy_last": {
                    "type": "boolean",
                    "description": "Copy the last row's text/formatting (default true). False blanks the new cells.",
                    "default": True,
                },
            },
            "required": ["file_path", "slide_locator", "table_locator"],
        },
    },
    {
        "name": "pptx_replace_picture",
        "description": "Replace a picture's image in place with a new image file, preserving the original position and size. The old image is re-pointed, never removed. picture_locator may be a 0-based picture index OR a shape name/text. Returns geometry in inches.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to the .pptx file",
                },
                "slide_locator": {
                    "description": "1-based slide index, or content string"
                },
                "picture_locator": {
                    "description": "0-based picture index (among pictures), or shape name/text"
                },
                "new_image_path": {
                    "type": "string",
                    "description": "Full path to the replacement image file",
                },
            },
            "required": [
                "file_path",
                "slide_locator",
                "picture_locator",
                "new_image_path",
            ],
        },
    },
    {
        "name": "pptx_add_autoshape",
        "description": "Add an autoshape (e.g. RECTANGLE, ROUNDED_RECTANGLE, OVAL) at the given geometry in inches, optionally filled. shape_type is an MSO_SHAPE name; unknown names fall back to RECTANGLE. Returns the created shape's geometry and fill read back from disk.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to the .pptx file",
                },
                "slide_locator": {
                    "description": "1-based slide index, or content string"
                },
                "shape_type": {
                    "type": "string",
                    "description": "MSO_SHAPE name, e.g. RECTANGLE, ROUNDED_RECTANGLE, OVAL",
                },
                "left": {"type": "number", "description": "Left position in inches"},
                "top": {"type": "number", "description": "Top position in inches"},
                "width": {"type": "number", "description": "Width in inches"},
                "height": {"type": "number", "description": "Height in inches"},
                "fill_hex": {
                    "type": "string",
                    "description": "Fill color, 6-digit hex no '#'. Optional.",
                },
            },
            "required": [
                "file_path",
                "slide_locator",
                "shape_type",
                "left",
                "top",
                "width",
                "height",
            ],
        },
    },
    {
        "name": "pptx_set_shape",
        "description": "Move/resize and/or recolor an existing shape. Any of left/top/width/height (inches) and fill_hex may be omitted to leave that property unchanged. shape_locator may be a 0-based shape index OR a shape name/text. Returns geometry and fill read back from disk.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to the .pptx file",
                },
                "slide_locator": {
                    "description": "1-based slide index, or content string"
                },
                "shape_locator": {
                    "description": "0-based shape index, or shape name/text"
                },
                "left": {
                    "type": "number",
                    "description": "Left in inches. Omit to leave unchanged.",
                },
                "top": {
                    "type": "number",
                    "description": "Top in inches. Omit to leave unchanged.",
                },
                "width": {
                    "type": "number",
                    "description": "Width in inches. Omit to leave unchanged.",
                },
                "height": {
                    "type": "number",
                    "description": "Height in inches. Omit to leave unchanged.",
                },
                "fill_hex": {
                    "type": "string",
                    "description": "Fill color, 6-digit hex no '#'. Omit to leave unchanged.",
                },
            },
            "required": ["file_path", "slide_locator", "shape_locator"],
        },
    },
    {
        "name": "pptx_add_hyperlink",
        "description": "Attach a hyperlink to the run within a shape's text that matches run_match (substring). Optionally recolors and underlines the run. shape_locator may be a 0-based shape index OR a shape name/text. Returns the resolved target read back from disk.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to the .pptx file",
                },
                "slide_locator": {
                    "description": "1-based slide index, or content string"
                },
                "shape_locator": {
                    "description": "0-based shape index, or shape name/text"
                },
                "run_match": {
                    "type": "string",
                    "description": "Substring identifying the text run to link",
                },
                "url": {"type": "string", "description": "The hyperlink target URL"},
                "color_hex": {
                    "type": "string",
                    "description": "Link font color, 6-digit hex no '#'. Default 1A73E8.",
                    "default": "1A73E8",
                },
                "underline": {
                    "type": "boolean",
                    "description": "Underline the linked run. Default true.",
                    "default": True,
                },
            },
            "required": [
                "file_path",
                "slide_locator",
                "shape_locator",
                "run_match",
                "url",
            ],
        },
    },
    {
        "name": "pptx_apply_theme",
        "description": (
            "Restyle/polish/reformat an EXISTING .pptx: apply background, font family, text colors, "
            "and table header/banding across a slide range in ONE call. USE THIS for 'polish this deck', "
            "'make it consistent', 'apply brand colors', 'reformat', 'dark theme', 'make a polished copy', "
            "'create a copy and format it', or any look-and-feel change to an existing deck. Workflow: "
            "copy the file, then call this on the copy. Do NOT rebuild with pptxgenjs (loses content, "
            "and the model narrates for 30K tokens then stops without acting). Do NOT write raw "
            "python-pptx in run_python (crashes on slide slicing prs.slides[:N] -> sldId.rId, "
            "RGBColor .red -> it's a tuple, missing Presentation() call). Decision: CREATING from "
            "scratch -> create_pptx; RESTYLING/copies of existing -> pptx_apply_theme; editing specific "
            "text -> update_pptx; editing a table cell -> pptx_write_table_cell. Every field optional: "
            "omit one to leave that aspect untouched. File saved in place; returns per-slide summary."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Full path to .pptx file, OneDrive item ID, or 'open' / 'open:Deck.pptx' for the active presentation via COM. OneDrive item IDs are automatically downloaded, edited, and re-uploaded.",
                },
                "slide_range": {
                    "type": "array",
                    "items": {"type": "integer"},
                    "description": "[start, end] 1-based inclusive slide numbers, e.g. [1, 32]. Omit for all slides.",
                },
                "bg_hex": {
                    "type": "string",
                    "description": "Solid background fill, 6-digit hex no '#', e.g. '0F1B2D'.",
                },
                "font_name": {
                    "type": "string",
                    "description": "Font family applied to every text run, e.g. 'Arial'.",
                },
                "title_color_hex": {
                    "type": "string",
                    "description": "Font color for title-placeholder runs, 6-digit hex no '#'.",
                },
                "body_color_hex": {
                    "type": "string",
                    "description": "Font color for all non-title runs, 6-digit hex no '#'.",
                },
                "table_header_fill_hex": {
                    "type": "string",
                    "description": "Solid fill for row 0 of every table, 6-digit hex no '#'.",
                },
                "table_band_fill_hex": {
                    "type": "string",
                    "description": "Solid fill for odd data rows (banding). Even rows stay transparent. 6-digit hex no '#'.",
                },
                "table_header_font_hex": {
                    "type": "string",
                    "description": "Font color for header-row cells, 6-digit hex no '#'.",
                },
                "table_body_font_hex": {
                    "type": "string",
                    "description": "Font color for data-row cells, 6-digit hex no '#'.",
                },
            },
            "required": ["file_path"],
        },
    },
]

TOOL_STATUS = {
    "get_pptx_info": "\U0001f4ca Inspecting PowerPoint...",
    "read_pptx": "\U0001f4ca Reading PowerPoint...",
    "create_pptx": "\U0001f4ca Creating PowerPoint...",
    "update_pptx": "\U0001f4ca Updating PowerPoint...",
    "pptx_list_shapes": "\U0001f4ca Listing slide shapes...",
    "pptx_read_table": "\U0001f4ca Reading table...",
    "pptx_write_table_cell": "\U0001f4ca Writing table cell...",
    "pptx_add_table_row": "\U0001f4ca Adding table row...",
    "pptx_replace_picture": "\U0001f4ca Replacing picture...",
    "pptx_add_autoshape": "\U0001f4ca Adding shape...",
    "pptx_set_shape": "\U0001f4ca Updating shape...",
    "pptx_add_hyperlink": "\U0001f4ca Adding hyperlink...",
    "pptx_apply_theme": "\U0001f4ca Applying theme...",
}


# ── OneDrive download / upload helper ────────────────────────────────────────


def _is_onedrive_ref(file_path: str) -> bool:
    """True if file_path is a OneDrive item ID or onedrive:// URI — not a local path."""
    if not file_path:
        return False
    if file_path.startswith("onedrive://"):
        return True
    # Graph item IDs are opaque alphanumeric strings, typically 20-40 chars,
    # with no path separators, dots, or spaces.
    if (
        20 <= len(file_path) <= 60
        and "\\" not in file_path
        and "/" not in file_path
        and "." not in file_path
        and " " not in file_path
        and not file_path.startswith("open")
    ):
        return True
    return False


def _parse_onedrive_ref(file_path: str) -> str:
    """Extract the Graph item ID from a onedrive:// URI or bare ID."""
    if file_path.startswith("onedrive://"):
        return file_path[len("onedrive://") :]
    return file_path


@contextlib.contextmanager
def _onedrive_pptx_context(item_id: str, drive_id: str = "", readonly: bool = False):
    """Download a .pptx from OneDrive to a temp file, yield (local_path, name),
    then upload the modified file back.  On readonly=True, skip the upload.

    Handles both personal OneDrive and SharePoint files. If drive_id is not
    provided, tries /me/drive/items/{id} first, then falls back to
    /me/drive/sharedWithMe to discover the correct drive_id for SharePoint items.
    """
    import httpx
    from pathlib import Path
    from .._m365.helpers import get_skill_client

    ONEDRIVE_SKILLS_DIR = Path(__file__).parent.parent / "m365-onedrive" / "scripts"
    gc = get_skill_client(ONEDRIVE_SKILLS_DIR)

    # Resolve metadata — use the shared lookup helper which handles
    # personal OneDrive + SharePoint (sharedWithMe fallback on 404).
    from skills.onedrive.tools import _try_direct_lookup_with_download_url

    direct = _try_direct_lookup_with_download_url(gc, item_id, drive_id)
    if not direct:
        raise RuntimeError(
            f"Could not resolve OneDrive item {item_id} (drive_id={drive_id!r}). "
            "The file may be on a SharePoint site — provide drive_id or use search_onedrive_files first."
        )
    meta = direct["meta"]
    name = meta.get("name", "presentation.pptx")
    direct_url = meta.get("@microsoft.graph.downloadUrl", "")
    real_id = meta.get("id", item_id)
    real_drive_id = direct.get("drive_id", drive_id)
    token = gc.get_token()

    # Download
    client = httpx.Client(timeout=httpx.Timeout(60.0), follow_redirects=True)
    try:
        if direct_url:
            r = client.get(direct_url)
        else:
            if real_drive_id:
                dl_url = f"https://graph.microsoft.com/v1.0/drives/{real_drive_id}/items/{real_id}/content"
            else:
                dl_url = (
                    f"https://graph.microsoft.com/v1.0/me/drive/items/{real_id}/content"
                )
            r = client.get(dl_url, headers={"Authorization": f"Bearer {token}"})
        r.raise_for_status()
        raw = r.content
    finally:
        client.close()

    if raw[:5] in (b"<!DOC", b"<html", b"<HTML"):
        raise RuntimeError(
            "OneDrive returned an HTML page instead of the file — token may have expired."
        )

    # Write to temp file
    suffix = f"_{name}" if name.endswith(".pptx") else "_presentation.pptx"
    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    tmp.write(raw)
    tmp.close()
    local_path = tmp.name

    try:
        yield local_path, name

        if not readonly:
            # Upload modified file back to OneDrive
            with open(local_path, "rb") as fh:
                updated = fh.read()

            # Use upload session — required for shared/co-authored files
            if real_drive_id:
                sess_url = f"https://graph.microsoft.com/v1.0/drives/{real_drive_id}/items/{real_id}/createUploadSession"
            else:
                sess_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{real_id}/createUploadSession"
            sess_client = httpx.Client(
                timeout=httpx.Timeout(30.0), follow_redirects=True
            )
            try:
                sess_resp = sess_client.post(
                    sess_url,
                    headers={
                        "Authorization": f"Bearer {token}",
                        "Content-Type": "application/json",
                    },
                    json={"item": {"@microsoft.graph.conflictBehavior": "replace"}},
                )
            finally:
                sess_client.close()

            if sess_resp.is_success:
                upload_url = sess_resp.json().get("uploadUrl", "")
                size = len(updated)
                up_client = httpx.Client(
                    timeout=httpx.Timeout(120.0), follow_redirects=True
                )
                try:
                    resp = up_client.put(
                        upload_url,
                        headers={
                            "Content-Range": f"bytes 0-{size - 1}/{size}",
                            "Content-Length": str(size),
                        },
                        content=updated,
                    )
                finally:
                    up_client.close()
            else:
                raise RuntimeError(
                    f"Upload session creation failed: {sess_resp.status_code} {sess_resp.text[:200]}"
                )
    finally:
        try:
            os.unlink(local_path)
        except OSError:
            pass


# ── Tool Handlers ────────────────────────────────────────────────────────────


def _tool_get_pptx_info(file_path: str) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=True
            ) as (local, name):
                result = _tool_get_pptx_info(local)
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download failed: {ex}"}
    try:
        if file_path.startswith("open"):
            from skills._office_com import get_ppt_app, get_ppt_presentation

            ppt, err = get_ppt_app()
            if err:
                return {"error": err}
            pres, err = get_ppt_presentation(ppt, file_path)
            if err:
                return {"error": err}
            slides_info = []
            for i in range(1, pres.Slides.Count + 1):
                slide = pres.Slides(i)
                title = ""
                try:
                    title = slide.Shapes.Title.TextFrame.TextRange.Text
                except Exception:
                    pass
                layout = slide.Layout
                slides_info.append({"number": i, "title": title, "layout_id": layout})
            return {
                "ok": True,
                "file": pres.FullName,
                "slide_count": pres.Slides.Count,
                "width": pres.PageSetup.SlideWidth,
                "height": pres.PageSetup.SlideHeight,
                "slides": slides_info,
            }
        else:
            from pptx import Presentation
            from pptx.util import Emu

            prs = Presentation(file_path)
            slides_info = []
            for idx, slide in enumerate(prs.slides, 1):
                title = slide.shapes.title.text if slide.shapes.title else ""
                layout = slide.slide_layout.name
                slides_info.append({"number": idx, "title": title, "layout": layout})
            return {
                "ok": True,
                "file": file_path,
                "slide_count": len(prs.slides),
                "width_inches": round(prs.slide_width / Emu(914400), 2),
                "height_inches": round(prs.slide_height / Emu(914400), 2),
                "slides": slides_info,
            }
    except Exception as e:
        return {"error": str(e)}


def _tool_read_pptx(file_path: str, slide_number: int = None) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=True
            ) as (local, name):
                result = _tool_read_pptx(local, slide_number)
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download failed: {ex}"}
    try:
        if file_path.startswith("open"):
            from skills._office_com import get_ppt_app, get_ppt_presentation

            ppt, err = get_ppt_app()
            if err:
                return {"error": err}
            pres, err = get_ppt_presentation(ppt, file_path)
            if err:
                return {"error": err}
            slides_data = []
            start = slide_number or 1
            end = slide_number or pres.Slides.Count
            for i in range(start, end + 1):
                slide = pres.Slides(i)
                title = ""
                try:
                    title = slide.Shapes.Title.TextFrame.TextRange.Text
                except Exception:
                    pass
                content = []
                for j in range(1, slide.Shapes.Count + 1):
                    shape = slide.Shapes(j)
                    if shape.HasTextFrame:
                        text = shape.TextFrame.TextRange.Text.strip()
                        if text and text != title:
                            content.append(text)
                notes = ""
                try:
                    notes = slide.NotesPage.Shapes(2).TextFrame.TextRange.Text.strip()
                except Exception:
                    pass
                slides_data.append(
                    {"number": i, "title": title, "content": content, "notes": notes}
                )
            return {"ok": True, "file": pres.FullName, "slides": slides_data}
        else:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides_data = []
            for idx, slide in enumerate(prs.slides, 1):
                if slide_number and idx != slide_number:
                    continue
                title = slide.shapes.title.text if slide.shapes.title else ""
                content = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text and text != title:
                            content.append(text)
                notes = ""
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    if notes_tf:
                        notes = notes_tf.text.strip()
                slides_data.append(
                    {"number": idx, "title": title, "content": content, "notes": notes}
                )
            return {"ok": True, "file": file_path, "slides": slides_data}
    except Exception as e:
        return {"error": str(e)}


def _tool_create_pptx(file_path: str, slides: list, author: str = "") -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()

        if author:
            prs.core_properties.author = author

        for slide_def in slides:
            layout_name = slide_def.get("layout", "title_content")
            layout_idx = LAYOUT_MAP.get(layout_name, 1)

            try:
                layout = prs.slide_layouts[layout_idx]
            except IndexError:
                layout = prs.slide_layouts[1]  # fallback to title+content

            slide = prs.slides.add_slide(layout)

            # Title
            title_text = slide_def.get("title", "")
            if title_text and slide.shapes.title:
                slide.shapes.title.text = title_text

            # Subtitle (title_slide only)
            subtitle_text = slide_def.get("subtitle", "")
            if subtitle_text and layout_name == "title_slide":
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        ph.text = subtitle_text
                        break

            # Content (bullets or text)
            content = slide_def.get("content")
            if content and layout_name not in ("blank", "title_only"):
                body_ph = None
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx not in (0,):  # skip title
                        body_ph = ph
                        break

                if body_ph and body_ph.has_text_frame:
                    tf = body_ph.text_frame
                    tf.clear()
                    items = content if isinstance(content, list) else [content]
                    for i, item in enumerate(items):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = item
                        p.font.size = Pt(18)

            # Speaker notes
            notes_text = slide_def.get("notes", "")
            if notes_text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text

            # Image
            image_def = slide_def.get("image")
            if image_def and os.path.exists(image_def.get("path", "")):
                left = Inches(image_def.get("left", 1))
                top = Inches(image_def.get("top", 2))
                width = Inches(image_def.get("width", 5))
                slide.shapes.add_picture(image_def["path"], left, top, width=width)

        prs.save(file_path)
        return {
            "ok": True,
            "message": f"Created {file_path} — {len(slides)} slide(s)",
            "file_path": file_path,
            "slide_count": len(slides),
        }
    except Exception as e:
        return {"error": str(e)}


def _tool_update_pptx(
    update_type: str,
    slide_number: int = None,
    new_text: str = "",
    file_path: str = "open",
    shape_index: int = 0,
    operations: list = None,
) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=False
            ) as (local, name):
                result = _tool_update_pptx(
                    update_type=update_type,
                    slide_number=slide_number,
                    new_text=new_text,
                    file_path=local,
                    shape_index=shape_index,
                    operations=operations,
                )
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download/upload failed: {ex}"}
    try:
        if update_type == "batch":
            if not operations:
                return {"error": "operations array is required for batch mode"}
            results = []
            for i, op in enumerate(operations):
                r = _tool_update_pptx(
                    file_path=file_path,
                    slide_number=op.get("slide_number"),
                    update_type=op.get("update_type", "title"),
                    new_text=op.get("new_text", ""),
                    shape_index=op.get("shape_index", 0),
                )
                results.append({"op": i + 1, **r})
                if r.get("error"):
                    break
            succeeded = sum(1 for r in results if r.get("ok"))
            return {
                "ok": True,
                "message": f"Batch: {succeeded}/{len(operations)} slides updated",
                "results": results,
            }

        if not slide_number:
            return {"error": "slide_number is required"}

        if file_path.startswith("open"):
            from skills._office_com import (
                get_ppt_app,
                get_ppt_presentation,
                save_com_document,
                get_file_info,
            )

            ppt, err = get_ppt_app()
            if err:
                return {"error": err}
            pres, err = get_ppt_presentation(ppt, file_path)
            if err:
                return {"error": err}
            finfo = get_file_info(pres, "ppt")
            slide = pres.Slides(slide_number)
            if update_type == "title":
                slide.Shapes.Title.TextFrame.TextRange.Text = new_text
            else:
                shape = slide.Shapes(shape_index + 1)
                shape.TextFrame.TextRange.Text = new_text
            save_com_document(pres, "ppt")
            return {
                "ok": True,
                **finfo,
                "message": f"Slide {slide_number} updated in {finfo['file_name']}",
            }
        else:
            from pptx import Presentation

            prs = Presentation(file_path)
            slide = prs.slides[slide_number - 1]
            if update_type == "title" and slide.shapes.title:
                slide.shapes.title.text = new_text
            elif update_type == "body":
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx != 0:
                        ph.text = new_text
                        break
            elif update_type == "shape":
                slide.shapes[shape_index].text_frame.text = new_text
            prs.save(file_path)
            return {"ok": True, "message": f"Saved {file_path}"}
    except Exception as e:
        return {"error": str(e)}


# ── Table/shape tools (python-pptx, closed-file only) ────────────────────────
# These fill the gap left by update_pptx (title/body/shape text only): cells by
# row/col, shape geometry/type, picture swaps, autoshapes, hyperlinks. Slides,
# tables and shapes are located by a content scan (or an integer index), never a
# hardcoded position. Every write re-opens the saved file to read back the
# mutated value, so success is proven from disk. No element is ever removed —
# pictures are re-pointed at a new image part, never deleted.

_EMU_PER_INCH = 914400


def _emu_to_in(v):
    if v is None:
        return None
    return round(v / _EMU_PER_INCH, 4)


def _norm_hex(h):
    if h is None:
        return None
    h = h.strip().lstrip("#").upper()
    return h or None


def _cell_styles(cell):
    """Extract {text, fill_hex, font_hex, font_bold} from a table cell, tolerating
    cells that have no explicit fill or run-level font (reports None, not crash)."""
    text = cell.text
    fill_hex = None
    try:
        fc = cell.fill.fore_color
        fill_hex = str(fc.rgb)
    except Exception:
        fill_hex = None
    font_hex = None
    font_bold = None
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            try:
                if r.font.color and r.font.color.type is not None:
                    font_hex = str(r.font.color.rgb)
            except Exception:
                pass
            if r.font.bold is not None:
                font_bold = r.font.bold
            if font_hex is not None or font_bold is not None:
                break
        if font_hex is not None or font_bold is not None:
            break
    return {
        "text": text,
        "fill_hex": font_or_none(fill_hex),
        "font_hex": font_or_none(font_hex),
        "font_bold": font_bold,
    }


def font_or_none(v):
    return v if v else None


def _shape_type_name(shape):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        st = shape.shape_type
    except Exception:
        return "UNKNOWN"
    mapping = {
        MSO_SHAPE_TYPE.TABLE: "TABLE",
        MSO_SHAPE_TYPE.PICTURE: "PICTURE",
        MSO_SHAPE_TYPE.TEXT_BOX: "TEXT_BOX",
        MSO_SHAPE_TYPE.AUTO_SHAPE: "AUTO_SHAPE",
        MSO_SHAPE_TYPE.GROUP: "GROUP",
    }
    return mapping.get(
        st, str(st).split(".")[-1].split(" ")[0] if st is not None else "UNKNOWN"
    )


def _shape_text(shape):
    try:
        if shape.has_text_frame:
            return shape.text_frame.text
    except Exception:
        pass
    return ""


def _resolve_slide(prs, locator):
    """1-based int / digit-string → slide by index. Otherwise scan slide title,
    shape text and table cell text for a case-insensitive substring match."""
    if isinstance(locator, int) or (
        isinstance(locator, str) and locator.strip().isdigit()
    ):
        idx = int(locator)
        if idx < 1 or idx > len(prs.slides):
            raise ValueError(f"slide index {idx} out of range (1..{len(prs.slides)})")
        return prs.slides[idx - 1]
    needle = str(locator).strip().lower()
    for slide in prs.slides:
        hay = []
        if slide.shapes.title:
            hay.append(slide.shapes.title.text)
        for shape in slide.shapes:
            hay.append(_shape_text(shape))
            try:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            hay.append(cell.text)
            except Exception:
                pass
        if any(needle in (t or "").lower() for t in hay):
            return slide
    raise ValueError(f"no slide matched locator {locator!r}")


def _resolve_table(slide, locator):
    """int → Nth table (0-based). Else match a header-row (row 0) cell text."""
    tables = [s for s in slide.shapes if getattr(s, "has_table", False)]
    if not tables:
        raise ValueError("slide has no tables")
    if isinstance(locator, int) or (
        isinstance(locator, str) and locator.strip().isdigit()
    ):
        idx = int(locator)
        if idx < 0 or idx >= len(tables):
            raise ValueError(f"table index {idx} out of range (0..{len(tables) - 1})")
        return tables[idx].table
    needle = str(locator).strip().lower()
    for gf in tables:
        tbl = gf.table
        for c in range(len(tbl.columns)):
            if needle in (tbl.cell(0, c).text or "").lower():
                return tbl
    raise ValueError(f"no table matched header locator {locator!r}")


def _resolve_shape(slide, locator, type_filter=None):
    """int → Nth shape (0-based, among type_filter if given). Else match shape
    name or text. type_filter is a set of type-name strings like {'PICTURE'}."""
    shapes = list(slide.shapes)
    if type_filter:
        shapes = [s for s in shapes if _shape_type_name(s) in type_filter]
    if not shapes:
        raise ValueError(f"slide has no shapes matching {type_filter}")
    if isinstance(locator, int) or (
        isinstance(locator, str) and locator.strip().isdigit()
    ):
        idx = int(locator)
        if idx < 0 or idx >= len(shapes):
            raise ValueError(f"shape index {idx} out of range (0..{len(shapes) - 1})")
        return shapes[idx]
    needle = str(locator).strip().lower()
    for s in shapes:
        if needle in (s.name or "").lower():
            return s
    for s in shapes:
        if needle in (_shape_text(s) or "").lower():
            return s
    raise ValueError(f"no shape matched locator {locator!r}")


def _slide_index(prs, slide):
    for i, s in enumerate(prs.slides):
        if s is slide:
            return i
    return None


def _tool_list_shapes(file_path: str, slide_locator) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=True
            ) as (local, name):
                result = _tool_list_shapes(local, slide_locator)
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download failed: {ex}"}
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slide = _resolve_slide(prs, slide_locator)
        out = []
        for i, shape in enumerate(slide.shapes):
            has_table = bool(getattr(shape, "has_table", False))
            preview = (_shape_text(shape) or "").strip().replace("\n", " ")
            if len(preview) > 80:
                preview = preview[:77] + "..."
            out.append(
                {
                    "index": i,
                    "name": shape.name,
                    "type": _shape_type_name(shape),
                    "left": _emu_to_in(shape.left),
                    "top": _emu_to_in(shape.top),
                    "width": _emu_to_in(shape.width),
                    "height": _emu_to_in(shape.height),
                    "has_table": has_table,
                    "text_preview": preview,
                }
            )
        return {"ok": True, "slide_index": _slide_index(prs, slide) + 1, "shapes": out}
    except Exception as e:
        return {"error": str(e)}


def _tool_read_table(file_path: str, slide_locator, table_locator) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=True
            ) as (local, name):
                result = _tool_read_table(local, slide_locator, table_locator)
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download failed: {ex}"}
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slide = _resolve_slide(prs, slide_locator)
        tbl = _resolve_table(slide, table_locator)
        rows = len(tbl.rows)
        cols = len(tbl.columns)
        grid = [
            [_cell_styles(tbl.cell(r, c)) for c in range(cols)] for r in range(rows)
        ]
        return {"ok": True, "rows": rows, "cols": cols, "grid": grid}
    except Exception as e:
        return {"error": str(e)}


def _apply_cell(cell, text, fill_hex, font_hex, bold):
    from pptx.dml.color import RGBColor

    if text is not None:
        cell.text = text
    if fill_hex is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(_norm_hex(fill_hex))
    if font_hex is not None or bold is not None:
        for p in cell.text_frame.paragraphs:
            runs = p.runs
            if not runs and (text is not None):
                continue
            for r in runs:
                if font_hex is not None:
                    r.font.color.rgb = RGBColor.from_string(_norm_hex(font_hex))
                if bold is not None:
                    r.font.bold = bold


def _tool_write_table_cell(
    file_path: str,
    slide_locator,
    table_locator,
    row: int,
    col: int,
    text: str = None,
    fill_hex: str = None,
    font_hex: str = None,
    bold: bool = None,
) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=False
            ) as (local, name):
                result = _tool_write_table_cell(
                    local,
                    slide_locator,
                    table_locator,
                    row,
                    col,
                    text,
                    fill_hex,
                    font_hex,
                    bold,
                )
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download/upload failed: {ex}"}
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slide = _resolve_slide(prs, slide_locator)
        tbl = _resolve_table(slide, table_locator)
        si = _slide_index(prs, slide)
        _apply_cell(tbl.cell(row, col), text, fill_hex, font_hex, bold)
        prs.save(file_path)
        # read back from disk
        prs2 = Presentation(file_path)
        slide2 = prs2.slides[si]
        tbl2 = _resolve_table(slide2, table_locator)
        after = _cell_styles(tbl2.cell(row, col))
        return {"ok": True, "row": row, "col": col, "cell_after": after}
    except Exception as e:
        return {"error": str(e)}


def _tool_add_table_row(
    file_path: str, slide_locator, table_locator, copy_last: bool = True
) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=False
            ) as (local, name):
                result = _tool_add_table_row(
                    local, slide_locator, table_locator, copy_last
                )
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download/upload failed: {ex}"}
    try:
        import copy as _copy
        from pptx import Presentation

        prs = Presentation(file_path)
        slide = _resolve_slide(prs, slide_locator)
        tbl = _resolve_table(slide, table_locator)
        si = _slide_index(prs, slide)
        tbl_el = tbl._tbl
        rows_el = tbl_el.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}tr"
        )
        last_tr = rows_el[-1]
        new_tr = _copy.deepcopy(last_tr)
        if not copy_last:
            a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for tc in new_tr.findall(f"{a}tc"):
                txbody = tc.find(f"{a}txBody")
                if txbody is not None:
                    for p in txbody.findall(f"{a}p"):
                        for run in p.findall(f"{a}r"):
                            t = run.find(f"{a}t")
                            if t is not None:
                                t.text = ""
        tbl_el.append(new_tr)
        prs.save(file_path)
        # read back
        prs2 = Presentation(file_path)
        tbl2 = _resolve_table(prs2.slides[si], table_locator)
        count = len(tbl2.rows)
        return {"ok": True, "new_row_index": count - 1, "row_count": count}
    except Exception as e:
        return {"error": str(e)}


def _tool_replace_picture(
    file_path: str, slide_locator, picture_locator, new_image_path: str
) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=False
            ) as (local, name):
                result = _tool_replace_picture(
                    local, slide_locator, picture_locator, new_image_path
                )
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download/upload failed: {ex}"}
    try:
        from pptx import Presentation

        if not os.path.exists(new_image_path):
            return {"error": f"image not found: {new_image_path}"}
        prs = Presentation(file_path)
        slide = _resolve_slide(prs, slide_locator)
        si = _slide_index(prs, slide)
        pic = _resolve_shape(slide, picture_locator, type_filter={"PICTURE"})
        # swap the image blob in place, preserving geometry — no .remove()
        image_part, rId = slide.part.get_or_add_image_part(new_image_path)
        blip = pic._element.blipFill.blip
        blip.rEmbed = rId
        left, top, width, height = pic.left, pic.top, pic.width, pic.height
        prs.save(file_path)
        return {
            "ok": True,
            "slide_index": si + 1,
            "left": _emu_to_in(left),
            "top": _emu_to_in(top),
            "width": _emu_to_in(width),
            "height": _emu_to_in(height),
        }
    except Exception as e:
        return {"error": str(e)}


_AUTOSHAPE_FALLBACK = "RECTANGLE"


def _tool_add_autoshape(
    file_path: str,
    slide_locator,
    shape_type: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_hex: str = None,
) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=False
            ) as (local, name):
                result = _tool_add_autoshape(
                    local, slide_locator, shape_type, left, top, width, height, fill_hex
                )
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download/upload failed: {ex}"}
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation(file_path)
        slide = _resolve_slide(prs, slide_locator)
        si = _slide_index(prs, slide)
        try:
            mso = getattr(MSO_SHAPE, str(shape_type).strip().upper())
        except AttributeError:
            mso = getattr(MSO_SHAPE, _AUTOSHAPE_FALLBACK)
        shp = slide.shapes.add_shape(
            mso, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if fill_hex is not None:
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor.from_string(_norm_hex(fill_hex))
        prs.save(file_path)
        # read back from disk
        prs2 = Presentation(file_path)
        slide2 = prs2.slides[si]
        new = list(slide2.shapes)[-1]
        after_fill = None
        try:
            after_fill = str(new.fill.fore_color.rgb)
        except Exception:
            pass
        return {
            "ok": True,
            "shape_index": len(slide2.shapes._spTree)
            and (len(list(slide2.shapes)) - 1),
            "name": new.name,
            "left": _emu_to_in(new.left),
            "top": _emu_to_in(new.top),
            "width": _emu_to_in(new.width),
            "height": _emu_to_in(new.height),
            "fill_hex": font_or_none(after_fill),
        }
    except Exception as e:
        return {"error": str(e)}


def _tool_set_shape(
    file_path: str,
    slide_locator,
    shape_locator,
    left: float = None,
    top: float = None,
    width: float = None,
    height: float = None,
    fill_hex: str = None,
) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=False
            ) as (local, name):
                result = _tool_set_shape(
                    local,
                    slide_locator,
                    shape_locator,
                    left,
                    top,
                    width,
                    height,
                    fill_hex,
                )
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download/upload failed: {ex}"}
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor

        prs = Presentation(file_path)
        slide = _resolve_slide(prs, slide_locator)
        si = _slide_index(prs, slide)
        shp = _resolve_shape(slide, shape_locator)
        name = shp.name
        if left is not None:
            shp.left = Inches(left)
        if top is not None:
            shp.top = Inches(top)
        if width is not None:
            shp.width = Inches(width)
        if height is not None:
            shp.height = Inches(height)
        if fill_hex is not None:
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor.from_string(_norm_hex(fill_hex))
        prs.save(file_path)
        # read back from disk by name
        prs2 = Presentation(file_path)
        slide2 = prs2.slides[si]
        shp2 = next((s for s in slide2.shapes if s.name == name), None)
        after_fill = None
        try:
            after_fill = str(shp2.fill.fore_color.rgb)
        except Exception:
            pass
        return {
            "ok": True,
            "name": name,
            "left": _emu_to_in(shp2.left),
            "top": _emu_to_in(shp2.top),
            "width": _emu_to_in(shp2.width),
            "height": _emu_to_in(shp2.height),
            "fill_hex": font_or_none(after_fill),
        }
    except Exception as e:
        return {"error": str(e)}


def _tool_add_hyperlink(
    file_path: str,
    slide_locator,
    shape_locator,
    run_match: str,
    url: str,
    color_hex: str = "1A73E8",
    underline: bool = True,
) -> dict:
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=False
            ) as (local, name):
                result = _tool_add_hyperlink(
                    local,
                    slide_locator,
                    shape_locator,
                    run_match,
                    url,
                    color_hex,
                    underline,
                )
                result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download/upload failed: {ex}"}
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        prs = Presentation(file_path)
        slide = _resolve_slide(prs, slide_locator)
        si = _slide_index(prs, slide)
        shp = _resolve_shape(slide, shape_locator)
        name = shp.name
        if not shp.has_text_frame:
            return {"error": "target shape has no text frame"}
        needle = str(run_match).strip().lower()
        target_run = None
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                if needle in (r.text or "").lower():
                    target_run = r
                    break
            if target_run is not None:
                break
        if target_run is None:
            return {"error": f"no run matched {run_match!r}"}
        target_run.hyperlink.address = url
        if color_hex is not None:
            target_run.font.color.rgb = RGBColor.from_string(_norm_hex(color_hex))
        if underline is not None:
            target_run.font.underline = underline
        prs.save(file_path)
        # read back from disk
        prs2 = Presentation(file_path)
        slide2 = prs2.slides[si]
        shp2 = next((s for s in slide2.shapes if s.name == name), None)
        resolved = None
        for p in shp2.text_frame.paragraphs:
            for r in p.runs:
                if r.hyperlink and r.hyperlink.address == url:
                    resolved = r.hyperlink.address
        return {"ok": True, "name": name, "target": resolved}
    except Exception as e:
        return {"error": str(e)}


# ── apply_theme ──────────────────────────────────────────────────────────────
#
# Bulk restyle across a slide range: one call replaces the per-shape raw
# python-pptx scripts the model kept getting wrong (slide slicing, RGBColor
# API, missing Presentation() call). Encapsulates every gotcha in tested code
# — the model calls one tool instead of authoring 60 lines of python-pptx.
#
# Design decisions:
#   * Iterates `for slide in prs.slides` (never slices — slicing a Slides
#     collection returns a list and breaks python-pptx's internal sldId lookup).
#   * RGBColor is a tuple subclass — built once via RGBColor.from_string, never
#     touched via .red/.green/.blue.
#   * All color inputs are 6-digit hex (no #), normalized via _norm_hex.
#   * Omit any theme field to leave that aspect untouched (partial update).
#   * Write reads back from disk and returns a per-slide summary so success is
#     proven, matching the convention of every other tool here.


def _set_slide_background(slide, rgb):
    """Set a slide's background to a solid fill. Overrides layout/master bg."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb


def _set_run_font(run, font_name, color_rgb=None, bold=None):
    """Apply font family + optional color/bold to a single run, tolerating
    runs that have no existing font object."""
    if font_name:
        run.font.name = font_name
    if color_rgb is not None:
        run.font.color.rgb = color_rgb
    if bold is not None:
        run.font.bold = bold


def _style_text_frames(slide, font_name, title_rgb, body_rgb):
    """Apply font family + color to every text run on the slide. Title
    placeholder runs get title_rgb; everything else gets body_rgb. Shapes
    inside groups are recursed so infographic slides are covered too."""

    def _is_title(shape):
        # `slide.shapes.title` returns a wrapper, not the same instance the
        # iteration yields, so identity comparison fails. Detect the title by
        # placeholder idx (0 == title) which is stable across instances.
        try:
            return shape.is_placeholder and shape.placeholder_format.idx == 0
        except Exception:
            return False

    def _walk(shapes):
        for shape in shapes:
            try:
                if shape.shape_type and hasattr(shape, "shapes"):
                    _walk(shape.shapes)
                    continue
            except Exception:
                pass
            if not getattr(shape, "has_text_frame", False):
                continue
            color = title_rgb if _is_title(shape) else body_rgb
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    _set_run_font(r, font_name, color_rgb=color)

    _walk(slide.shapes)


def _style_tables(
    slide, header_fill_rgb, band_fill_rgb, header_font_rgb, body_font_rgb
):
    """Apply header fill + banded data rows to every table on the slide.
    Header = row 0; data rows alternate band_fill_rgb / transparent."""
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        tbl = shape.table
        for r in range(len(tbl.rows)):
            for c in range(len(tbl.columns)):
                cell = tbl.cell(r, c)
                if r == 0 and header_fill_rgb is not None:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_fill_rgb
                elif r > 0 and band_fill_rgb is not None and r % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = band_fill_rgb
                # font color: header rows get header_font_rgb, else body_font_rgb
                font_rgb = header_font_rgb if r == 0 else body_font_rgb
                if font_rgb is None:
                    continue
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = font_rgb


def _tool_apply_theme(
    file_path: str,
    slide_range: list = None,
    bg_hex: str = None,
    font_name: str = None,
    title_color_hex: str = None,
    body_color_hex: str = None,
    table_header_fill_hex: str = None,
    table_band_fill_hex: str = None,
    table_header_font_hex: str = None,
    table_body_font_hex: str = None,
) -> dict:
    """Apply a unified visual theme across a range of slides in one call.

    This is the right tool for "make this deck consistent" / "polish the look"
    / "apply our brand colors across all slides" — it replaces the per-shape
    raw python-pptx scripts that tend to get the slide-iteration / RGBColor
    API wrong. Every field is optional: omit one to leave that aspect untouched.

    Args:
        file_path: .pptx path, OneDrive item ID, or 'open'/'open:Name' for COM.
        slide_range: [start, end] 1-based inclusive slide numbers, e.g. [1, 32].
            None (default) = all slides.
        bg_hex: solid background fill, 6-digit hex no '#', e.g. "0F1B2D".
        font_name: font family applied to every text run, e.g. "Arial".
        title_color_hex: font color for title-placeholder runs.
        body_color_hex: font color for all non-title runs.
        table_header_fill_hex: solid fill for row 0 of every table.
        table_band_fill_hex: solid fill for odd data rows (banding). Even data
            rows are left transparent so banding reads against the bg.
        table_header_font_hex: font color for header-row cells.
        table_body_font_hex: font color for data-row cells.

    Returns:
        {"ok": True, "slides_changed": N, "details": [{slide, bg, font, tables}, ...]}
        on success; {"error": str} on failure. The file is saved in place.
    """
    if _is_onedrive_ref(file_path):
        try:
            with _onedrive_pptx_context(
                _parse_onedrive_ref(file_path), readonly=False
            ) as (local, name):
                result = _tool_apply_theme(
                    local,
                    slide_range,
                    bg_hex,
                    font_name,
                    title_color_hex,
                    body_color_hex,
                    table_header_fill_hex,
                    table_band_fill_hex,
                    table_header_font_hex,
                    table_body_font_hex,
                )
                if result.get("ok"):
                    result["file"] = f"OneDrive: {name}"
                return result
        except Exception as ex:
            return {"error": f"OneDrive download/upload failed: {ex}"}
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        def _maybe_hex(h):
            h = _norm_hex(h)
            return RGBColor.from_string(h) if h else None

        bg_rgb = _maybe_hex(bg_hex)
        title_rgb = _maybe_hex(title_color_hex)
        body_rgb = _maybe_hex(body_color_hex)
        th_fill = _maybe_hex(table_header_fill_hex)
        tb_fill = _maybe_hex(table_band_fill_hex)
        th_font = _maybe_hex(table_header_font_hex)
        tb_font = _maybe_hex(table_body_font_hex)

        prs = Presentation(file_path)
        total = len(prs.slides)
        if slide_range is None:
            lo, hi = 1, total
        else:
            lo, hi = int(slide_range[0]), int(slide_range[1])
            if lo < 1 or hi > total or lo > hi:
                raise ValueError(
                    f"slide_range {slide_range} out of bounds (1..{total})"
                )

        details = []
        changed = 0
        # Iterate by enumeration — NEVER slice prs.slides (slicing returns a
        # list and breaks python-pptx's internal sldId.rId lookup).
        for idx, slide in enumerate(prs.slides, start=1):
            if idx < lo or idx > hi:
                continue
            had_table = False
            if bg_rgb is not None:
                _set_slide_background(slide, bg_rgb)
            if font_name or title_rgb is not None or body_rgb is not None:
                _style_text_frames(
                    slide,
                    font_name or None,
                    title_rgb,
                    body_rgb,
                )
            table_count = sum(1 for s in slide.shapes if getattr(s, "has_table", False))
            if table_count and (
                th_fill is not None
                or tb_fill is not None
                or th_font is not None
                or tb_font is not None
            ):
                _style_tables(slide, th_fill, tb_fill, th_font, tb_font)
                had_table = True
            changed += 1
            details.append(
                {
                    "slide": idx,
                    "bg": _norm_hex(bg_hex),
                    "font": font_name or None,
                    "tables": table_count if had_table or table_count else 0,
                }
            )

        prs.save(file_path)
        return {
            "ok": True,
            "slides_changed": changed,
            "slide_range": [lo, hi],
            "details": details,
        }
    except Exception as e:
        return {"error": str(e)}


TOOL_HANDLERS = {
    "get_pptx_info": _tool_get_pptx_info,
    "read_pptx": _tool_read_pptx,
    "create_pptx": _tool_create_pptx,
    "update_pptx": _tool_update_pptx,
    "pptx_list_shapes": _tool_list_shapes,
    "pptx_read_table": _tool_read_table,
    "pptx_write_table_cell": _tool_write_table_cell,
    "pptx_add_table_row": _tool_add_table_row,
    "pptx_replace_picture": _tool_replace_picture,
    "pptx_add_autoshape": _tool_add_autoshape,
    "pptx_set_shape": _tool_set_shape,
    "pptx_add_hyperlink": _tool_add_hyperlink,
    "pptx_apply_theme": _tool_apply_theme,
}
