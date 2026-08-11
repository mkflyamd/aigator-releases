"""read_onedrive_file must fall back to sharedWithMe when a direct item lookup
fails (400 OR 404) for a SharePoint file pinned WITHOUT a drive_id.

Native-pane OneDrive pins often carry a file_id but no drive_id. Calling
/me/drive/items/<id> then fails — 404 if the id is personal-OneDrive-shaped but
not found, or 400 ("Invalid request") if the id is a SharePoint base64url id
that Graph rejects for that endpoint. The fix: on 400/404 (no drive_id), scan
/me/drive/sharedWithMe to discover the correct driveId, then retry via
/drives/{driveId}/items/{itemId}. If that fails too, fall back to name-search.
"""
import io
import zipfile
from unittest.mock import MagicMock, patch

import pytest


def _make_pptx_bytes() -> bytes:
    """Minimal valid .pptx (a zip with the right content types) is overkill;
    python-pptx needs a real package, so we build a tiny one via python-pptx."""
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class _GraphError(Exception):
    def __init__(self, msg, status_code):
        super().__init__(msg)
        self.status_code = status_code


def _NotFound():
    return _GraphError("Graph API 404: The resource could not be found.", 404)


def _BadRequest():
    return _GraphError("Graph API 400: Invalid request.", 400)


def test_shared_fallback_resolves_drive_id_then_reads(monkeypatch):
    from skills.onedrive import tools

    raw = _make_pptx_bytes()

    gc = MagicMock()

    # 1st call: /me/drive/items/<id> -> 404 (file not on personal drive)
    # 2nd call: /me/drive/sharedWithMe -> lists the item with its real driveId
    # 3rd call: /drives/{driveId}/items/{id} -> metadata with downloadUrl
    def _get(path, params=None, **kw):
        if path == "/me/drive/items/FILE1":
            raise _NotFound()
        if path == "/me/drive/sharedWithMe":
            return {"value": [{
                "id": "FILE1",
                "name": "Deck.pptx",
                "remoteItem": {
                    "id": "FILE1",
                    "parentReference": {"driveId": "DRIVE_SP"},
                },
            }]}
        if path == "/drives/DRIVE_SP/items/FILE1":
            return {
                "id": "FILE1", "name": "Deck.pptx", "size": len(raw),
                "webUrl": "https://amd.sharepoint.com/sites/x/Deck.pptx",
                "@microsoft.graph.downloadUrl": "https://dl.example/Deck.pptx",
            }
        raise AssertionError(f"unexpected Graph path: {path}")

    gc.get.side_effect = _get
    gc.get_token.return_value = "TOKEN"

    # Stub the HTTP download pool to return our pptx bytes.
    resp = MagicMock()
    resp.content = raw
    resp.raise_for_status.return_value = None
    pool = MagicMock()
    pool.get.return_value = resp
    pool.is_closed = False

    with patch("skills._m365.helpers.get_skill_client", return_value=gc), \
         patch("httpx.Client", return_value=pool):
        # ensure a fresh pool is created
        if hasattr(tools._tool_read_onedrive_file, "_pool"):
            del tools._tool_read_onedrive_file._pool
        result = tools._tool_read_onedrive_file(file_id="FILE1")

    assert "error" not in result, result
    assert result["name"] == "Deck.pptx"
    # Proves the sharedWithMe fallback was exercised.
    paths = [c.args[0] for c in gc.get.call_args_list]
    assert "/me/drive/sharedWithMe" in paths
    assert "/drives/DRIVE_SP/items/FILE1" in paths


def test_400_triggers_shared_fallback_for_sharepoint_id(monkeypatch):
    """The actual bug: a SharePoint base64url id used against /me/drive/items
    returns 400 (not 404). The read path must still fall back to sharedWithMe."""
    from skills.onedrive import tools

    raw = _make_pptx_bytes()
    sp_id = "BcxSFOiEokG2yCDVW27EbPqkJCY-BdVHuh92nXumH2KnqVx5h7oRToXNciyM-LKPGEyTkbthQUWfLaBLAPiXdw"

    gc = MagicMock()

    def _get(path, params=None, **kw):
        if path == f"/me/drive/items/{sp_id}":
            raise _BadRequest()  # 400, not 404 — the bug
        if path == "/me/drive/sharedWithMe":
            return {"value": [{
                "id": sp_id, "name": "Planning.pptx",
                "remoteItem": {"id": sp_id,
                               "parentReference": {"driveId": "DRV_SP"}},
            }]}
        if path == f"/drives/DRV_SP/items/{sp_id}":
            return {
                "id": sp_id, "name": "Planning.pptx", "size": len(raw),
                "webUrl": "https://amd.sharepoint.com/x/Planning.pptx",
                "@microsoft.graph.downloadUrl": "https://dl.example/Planning.pptx",
            }
        raise AssertionError(f"unexpected Graph path: {path}")

    gc.get.side_effect = _get
    gc.get_token.return_value = "TOKEN"

    resp = MagicMock()
    resp.content = raw
    resp.raise_for_status.return_value = None
    pool = MagicMock()
    pool.get.return_value = resp
    pool.is_closed = False

    with patch("skills._m365.helpers.get_skill_client", return_value=gc), \
         patch("httpx.Client", return_value=pool):
        if hasattr(tools._tool_read_onedrive_file, "_pool"):
            del tools._tool_read_onedrive_file._pool
        result = tools._tool_read_onedrive_file(file_id=sp_id)

    assert "error" not in result, result
    assert result["name"] == "Planning.pptx"
    paths = [c.args[0] for c in gc.get.call_args_list]
    assert "/me/drive/sharedWithMe" in paths
    assert f"/drives/DRV_SP/items/{sp_id}" in paths


def test_404_without_shared_match_returns_error(monkeypatch):
    """When direct lookup 404s, sharedWithMe has no match, and there's no
    filename to fall back to name-search, the read returns a clear error dict
    instead of raising. The caller (agent) gets a actionable message."""
    from skills.onedrive import tools

    gc = MagicMock()

    def _get(path, params=None, **kw):
        if path == "/me/drive/items/FILE_X":
            raise _NotFound()
        if path == "/me/drive/sharedWithMe":
            return {"value": []}  # no match
        raise AssertionError(f"unexpected Graph path: {path}")

    gc.get.side_effect = _get

    with patch("skills._m365.helpers.get_skill_client", return_value=gc):
        result = tools._tool_read_onedrive_file(file_id="FILE_X")

    assert "error" in result
    assert "FILE_X" in result["error"]
    # Must have tried sharedWithMe before giving up.
    paths = [c.args[0] for c in gc.get.call_args_list]
    assert "/me/drive/sharedWithMe" in paths
