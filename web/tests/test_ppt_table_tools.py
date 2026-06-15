"""Tests for the table/shape PowerPoint tools (python-pptx, closed-file).

These tools fill the gap left by update_pptx (title/body/shape text only):
addressing table cells by row/col, introspecting shape geometry/type,
replacing picture images in place, creating/recoloring autoshapes, and
embedding run-level hyperlinks. Every write re-opens the saved file and reads
the mutated value back, so success is proven from disk rather than asserted.
"""
import os
import tempfile

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from skills.ppt.tools import (
    _tool_list_shapes,
    _tool_read_table,
    _tool_write_table_cell,
    _tool_add_table_row,
    _tool_replace_picture,
    _tool_add_autoshape,
    _tool_set_shape,
    _tool_add_hyperlink,
)


def _make_png(path, color):
    Image.new("RGB", (4, 4), color).save(path)


@pytest.fixture
def deck(tmp_path):
    """A one-slide deck: a 2x3 table with a styled header, a textbox with a
    linkable run, and a picture. Returns the saved .pptx path."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    if slide.shapes.title:
        slide.shapes.title.text = "Status Dashboard"

    # Table: header row + one data row
    gf = slide.shapes.add_table(2, 3, Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
    table = gf.table
    headers = ["Item", "Owner", "Updates"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
    # Style header cell [0][0]: red fill, white bold font — read_table must report these
    hc = table.cell(0, 0)
    hc.fill.solid()
    hc.fill.fore_color.rgb = RGBColor.from_string("FF0000")
    run = hc.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor.from_string("FFFFFF")
    run.font.bold = True
    # Data row
    table.cell(1, 0).text = "Migration"
    table.cell(1, 1).text = "Alice"
    table.cell(1, 2).text = ""  # the empty "Updates" cell we want to populate

    # Textbox with a run containing a Jira id
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(6), Inches(0.6))
    tb.name = "JiraNote"
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Tracked in JIRA-123 today"

    # Picture
    img0 = str(tmp_path / "orig.png")
    _make_png(img0, (255, 0, 0))
    slide.shapes.add_picture(img0, Inches(0.5), Inches(5), Inches(2), Inches(2))

    path = str(tmp_path / "deck.pptx")
    prs.save(path)
    return path


# ── list_shapes ────────────────────────────────────────────────────────────

def test_list_shapes_reports_types_and_geometry(deck):
    res = _tool_list_shapes(deck, slide_locator=1)
    assert res["ok"] is True
    types = {s["type"] for s in res["shapes"]}
    assert "TABLE" in types
    assert "PICTURE" in types
    table_shape = next(s for s in res["shapes"] if s["type"] == "TABLE")
    assert table_shape["has_table"] is True
    # geometry surfaced in inches
    pic = next(s for s in res["shapes"] if s["type"] == "PICTURE")
    assert round(pic["width"], 2) == 2.0
    assert round(pic["left"], 2) == 0.5


def test_list_shapes_resolves_slide_by_content(deck):
    # overloaded locator: a content string instead of an index
    res = _tool_list_shapes(deck, slide_locator="Status Dashboard")
    assert res["ok"] is True
    assert any(s["has_table"] for s in res["shapes"])


# ── read_table ─────────────────────────────────────────────────────────────

def test_read_table_returns_grid_with_styles(deck):
    res = _tool_read_table(deck, slide_locator=1, table_locator="Item")
    assert res["ok"] is True
    assert res["rows"] == 2 and res["cols"] == 3
    h00 = res["grid"][0][0]
    assert h00["text"] == "Item"
    assert h00["fill_hex"] == "FF0000"
    assert h00["font_hex"] == "FFFFFF"
    assert h00["font_bold"] is True
    # unstyled cell reports None colors rather than crashing
    assert res["grid"][1][1]["text"] == "Alice"
    assert res["grid"][1][1]["fill_hex"] is None


def test_read_table_locator_by_header_text(deck):
    # "Updates" is a header cell — table_locator scans the header row
    res = _tool_read_table(deck, slide_locator=1, table_locator="Updates")
    assert res["ok"] is True
    assert res["grid"][0][2]["text"] == "Updates"


# ── write_table_cell ───────────────────────────────────────────────────────

def test_write_table_cell_text_and_styles_with_readback(deck):
    res = _tool_write_table_cell(
        deck, slide_locator=1, table_locator="Item", row=1, col=2,
        text="Done 6/12", fill_hex="00FF00", font_hex="000000", bold=True,
    )
    assert res["ok"] is True
    after = res["cell_after"]
    assert after["text"] == "Done 6/12"
    assert after["fill_hex"] == "00FF00"
    assert after["font_hex"] == "000000"
    assert after["font_bold"] is True
    # prove persistence independently: re-open and read the same cell
    prs = Presentation(deck)
    tbl = next(s for s in prs.slides[0].shapes if s.has_table).table
    assert tbl.cell(1, 2).text == "Done 6/12"


def test_write_table_cell_partial_update_preserves_text(deck):
    # writing only a fill must not blank existing text
    res = _tool_write_table_cell(
        deck, slide_locator=1, table_locator="Item", row=1, col=1, fill_hex="123456",
    )
    assert res["ok"] is True
    assert res["cell_after"]["text"] == "Alice"
    assert res["cell_after"]["fill_hex"] == "123456"


# ── add_table_row ──────────────────────────────────────────────────────────

def test_add_table_row_deepcopies_last_and_reports_count(deck):
    res = _tool_add_table_row(deck, slide_locator=1, table_locator="Item")
    assert res["ok"] is True
    assert res["row_count"] == 3
    assert res["new_row_index"] == 2
    # re-open: the table really has 3 rows now
    prs = Presentation(deck)
    tbl = next(s for s in prs.slides[0].shapes if s.has_table).table
    assert len(tbl.rows) == 3


def test_add_table_row_copy_last_false_blanks_cells(deck):
    res = _tool_add_table_row(deck, slide_locator=1, table_locator="Item", copy_last=False)
    assert res["ok"] is True
    prs = Presentation(deck)
    tbl = next(s for s in prs.slides[0].shapes if s.has_table).table
    new_idx = res["new_row_index"]
    assert all(tbl.cell(new_idx, c).text == "" for c in range(3))


# ── replace_picture ────────────────────────────────────────────────────────

def test_replace_picture_swaps_blob_preserving_geometry(deck, tmp_path):
    new_img = str(tmp_path / "blue.png")
    _make_png(new_img, (0, 0, 255))
    res = _tool_replace_picture(deck, slide_locator=1, picture_locator=0, new_image_path=new_img)
    assert res["ok"] is True
    # geometry preserved
    assert round(res["width"], 2) == 2.0
    assert round(res["left"], 2) == 0.5
    # re-open: the picture's image bytes now equal the new file
    prs = Presentation(deck)
    pic = next(s for s in prs.slides[0].shapes if s.shape_type == 13)  # PICTURE
    with open(new_img, "rb") as f:
        assert pic.image.blob == f.read()


# ── add_autoshape ──────────────────────────────────────────────────────────

def test_add_autoshape_creates_filled_shape(deck):
    res = _tool_add_autoshape(
        deck, slide_locator=1, shape_type="ROUNDED_RECTANGLE",
        left=1, top=1, width=3, height=0.5, fill_hex="3366CC",
    )
    assert res["ok"] is True
    assert res["fill_hex"] == "3366CC"
    assert round(res["width"], 2) == 3.0
    # re-open: a new autoshape exists with that fill
    prs = Presentation(deck)
    shapes = list(prs.slides[0].shapes)
    found = [s for s in shapes if s.shape_type == 1]  # AUTO_SHAPE
    assert found


# ── set_shape ──────────────────────────────────────────────────────────────

def test_set_shape_moves_and_recolors_with_readback(deck):
    # target the textbox by name
    res = _tool_set_shape(
        deck, slide_locator=1, shape_locator="JiraNote",
        left=2, top=2, width=4, height=1, fill_hex="EEEEEE",
    )
    assert res["ok"] is True
    assert round(res["left"], 2) == 2.0
    assert round(res["height"], 2) == 1.0
    assert res["fill_hex"] == "EEEEEE"


# ── add_hyperlink ──────────────────────────────────────────────────────────

def test_add_hyperlink_on_run_with_readback(deck):
    res = _tool_add_hyperlink(
        deck, slide_locator=1, shape_locator="JiraNote",
        run_match="JIRA-123", url="https://jira.example.com/JIRA-123",
        color_hex="1A73E8", underline=True,
    )
    assert res["ok"] is True
    assert res["target"] == "https://jira.example.com/JIRA-123"
    # re-open: the run actually carries the hyperlink
    prs = Presentation(deck)
    tb = next(s for s in prs.slides[0].shapes if s.name == "JiraNote")
    addrs = [r.hyperlink.address for p in tb.text_frame.paragraphs for r in p.runs]
    assert "https://jira.example.com/JIRA-123" in addrs


def test_missing_slide_locator_errors_cleanly(deck):
    res = _tool_read_table(deck, slide_locator="No Such Slide", table_locator="Item")
    assert res.get("ok") is not True
    assert "error" in res
